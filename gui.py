"""
A tkinter-based GUI for the Dynamic Ollama Assistant.

Conversation state is persisted to a local file `conversation_state.json` in
the project root. This file is auto-created by the app on first save/close
and is ignored by Git by default, so it stays on your machine.
"""

import contextlib
import datetime
import json
import logging
import os
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, simpledialog, ttk
from tkinter import font as tkfont

import requests
import pandas as pd

from dynamic_ollama_assistant import (
    CSV_DIR,
    CSV_GLOB,
    EXCEL_GLOB,
    DEFAULT_MODEL as OLLAMA_MODEL,
    load_prompt_catalog,
    query_ollama_chat_for_gui,
    build_system_prompt,
    find_placeholders,
)
from web_scraper import scrape_web_content, crawl_website
from authenticated_scraper import (
    scrape_with_login_sync,
    analyze_login_form_sync,
    navigate_and_scrape_sync,
    playwright_crawl_sync,
)
from auth_dialogs import (
    ManualSelectorDialog,
    VerificationRequiredDialog,
    LoginAnalysisDialog,
    NavigationDialog,
)
from file_utils import process_uploaded_file, validate_url, aggregate_parsed_content
from docling.document_converter import DocumentConverter


class Tooltip:
    """A very small tooltip that shows on hover with dynamic text."""

    def __init__(self, widget, text_provider):
        self.widget = widget
        self.text_provider = text_provider  # callable returning str
        self.tipwindow = None
        self._after_id = None
        widget.bind("<Enter>", self._on_enter)
        widget.bind("<Leave>", self._on_leave)
        widget.bind("<ButtonPress>", self._on_leave)

    def _on_enter(self, _event=None):
        # small delay before showing
        self._schedule(400, self._show_tooltip)

    def _on_leave(self, _event=None):
        self._cancel()
        self._hide_tooltip()

    def _schedule(self, delay_ms, func):
        self._cancel()
        self._after_id = self.widget.after(delay_ms, func)

    def _cancel(self):
        if self._after_id is not None:
            self.widget.after_cancel(self._after_id)
            self._after_id = None

    def _show_tooltip(self):
        if self.tipwindow is not None:
            return
        text = (
            self.text_provider()
            if callable(self.text_provider)
            else str(self.text_provider)
        )
        if not text:
            return
        x = self.widget.winfo_rootx() + 10
        y = self.widget.winfo_rooty() + self.widget.winfo_height() + 8
        tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        label = tk.Label(
            tw,
            text=text,
            justify=tk.LEFT,
            background="#ffffe0",
            relief=tk.SOLID,
            borderwidth=1,
            padx=6,
            pady=4,
            wraplength=400,
        )
        label.pack(ipadx=1)
        self.tipwindow = tw

    def _hide_tooltip(self):
        if self.tipwindow is not None:
            self.tipwindow.destroy()
            self.tipwindow = None


class UIComponents:
    """
    A container for all UI widgets of the OllamaGUI.

    Fields:
        - parent: The OllamaGUI instance.
        :type parent: OllamaGUI

        - placeholder_entries: A dictionary of placeholder entries.
        :type placeholder_entries: dict

        - search_var: The search variable.
        :type search_var: tk.StringVar

        - prompt_tree: The treeview for displaying prompts.
        :type prompt_tree: ttk.Treeview

        - selected_prompt_label: The label for the selected prompt.
        :type selected_prompt_label: ttk.Label

        - placeholder_frame: The frame for placeholder entries.
        :type placeholder_frame: ttk.Frame

        - description_frame: The frame for prompt descriptions.
        :type description_frame: ttk.Frame

        - chat_history: The chat history widget.
        :type chat_history: scrolledtext.ScrolledText
    """

    def __init__(self, parent):
        """Initialize and create all UI components."""
        self.parent = parent  # The parent is the OllamaGUI instance
        self.placeholder_entries = {}

        # Authentication state tracking
        self.login_analyzed = False
        self.login_selectors = None
        self.authenticated_session = False

        # --- Main layout frames ---
        main_frame = ttk.Frame(parent, padding="10")
        main_frame.pack(fill="both", expand=True)

        paned_window = ttk.PanedWindow(main_frame, orient="horizontal")
        paned_window.pack(fill="both", expand=True)

        sidebar = ttk.Frame(paned_window, width=400)
        paned_window.add(sidebar, weight=1)

        content_area = ttk.Frame(paned_window)
        paned_window.add(content_area, weight=3)

        # --- Sidebar components ---
        self._create_sidebar_widgets(sidebar)

        # --- Content area components ---
        self._create_content_area_widgets(content_area)

    def _create_sidebar_widgets(self, sidebar):
        """Create widgets for the sidebar."""
        search_frame = ttk.Frame(sidebar)
        search_frame.pack(fill="x", pady=5)
        self.search_var = tk.StringVar()
        search_entry = ttk.Entry(search_frame, textvariable=self.search_var)
        search_entry.pack(fill="x", padx=5, ipady=2)
        search_entry.insert(0, "Search prompts...")
        search_entry.bind("<FocusIn>", self.parent.clear_placeholder)
        search_entry.bind("<FocusOut>", self.parent.add_placeholder)
        self.search_var.trace_add("write", self.parent.on_search)

        button_frame = ttk.Frame(sidebar)
        button_frame.pack(fill="x", pady=5)
        expand_button = ttk.Button(
            button_frame, text="Expand All", command=self.parent.expand_all
        )
        expand_button.pack(side="left", padx=5, expand=True, fill="x")
        collapse_button = ttk.Button(
            button_frame, text="Collapse All", command=self.parent.collapse_all
        )
        collapse_button.pack(side="left", padx=5, expand=True, fill="x")

        self.prompt_tree = ttk.Treeview(sidebar)
        self.prompt_tree.pack(fill="both", expand=True)
        self.prompt_tree.bind("<<TreeviewSelect>>", self.parent.on_prompt_select)

    def update_response(self, chunk):
        """Update the chat display with a chunk of the response."""
        self.chat_history.config(state=tk.NORMAL)
        self.chat_history.insert(tk.END, chunk)
        self.chat_history.see(tk.END)
        self.chat_history.config(state=tk.DISABLED)
        self.chat_history.update_idletasks()

    def _create_content_area_widgets(self, content_area):
        """Create widgets for the main content area."""
        prompt_frame = ttk.LabelFrame(content_area, text="Prompt Details")
        prompt_frame.pack(fill="x", pady=(0, 10), anchor="n")

        self.selected_prompt_label = ttk.Label(
            prompt_frame,
            text="No Prompt Selected",
            font=("Proxima Nova Alt", 10, "italic"),
        )
        self.selected_prompt_label.pack(pady=5)

        self.description_frame = ttk.Frame(prompt_frame, padding="5 0 0 0")
        self.description_frame.pack(fill="x", expand=True, padx=10, pady=(0, 10))

        self.placeholder_frame = ttk.Frame(prompt_frame)
        self.placeholder_frame.pack(fill="x", padx=10, pady=10)

        file_ops_frame = ttk.LabelFrame(content_area, text="File Operations")
        file_ops_frame.pack(fill="x", pady=5, anchor="n")

        # Top row: Upload Files button
        top_row = ttk.Frame(file_ops_frame)
        top_row.pack(fill="x", padx=5, pady=5)

        self.upload_multi_button = ttk.Button(
            top_row,
            text="Upload Files",
            command=self.parent.upload_and_parse_files,
        )
        self.upload_multi_button.pack(side="left")

        # Bottom row: Web scraping controls
        web_row = ttk.Frame(file_ops_frame)
        web_row.pack(fill="x", padx=5, pady=(0, 5))

        ttk.Label(web_row, text="URL:").pack(side="left")

        self.url_entry = ttk.Entry(web_row, width=40)
        self.url_entry.pack(side="left", padx=(5, 5), fill="x", expand=True)

        self.remote_endpoint_var = tk.StringVar(value=os.getenv("PLAYWRIGHT_REMOTE_ENDPOINT", ""))
        self.endpoint_entry = ttk.Entry(web_row, width=38, textvariable=self.remote_endpoint_var)
        self.endpoint_entry.pack(side="left", padx=(5, 5), fill="x", expand=True)
        self.endpoint_entry.configure(
            foreground="white" if self.remote_endpoint_var.get() else "gray"
        )

        self.remote_endpoint_var.trace_add("write", self._refresh_endpoint_entry_color)

        # Set up placeholder behavior for URL field
        self.url_placeholder = "Enter URL to scrape..."
        self.url_entry.insert(0, self.url_placeholder)
        self.url_entry.config(foreground="gray")

        def url_focus_in(event):
            if self.url_entry.get() == self.url_placeholder:
                self.url_entry.delete(0, tk.END)
                self.url_entry.config(foreground="white")

        def url_focus_out(event):
            if not self.url_entry.get().strip():
                self.url_entry.insert(0, self.url_placeholder)
                self.url_entry.config(foreground="gray")

        self.url_entry.bind("<FocusIn>", url_focus_in)
        self.url_entry.bind("<FocusOut>", url_focus_out)

        self.scrape_button = ttk.Button(
            web_row,
            text="Scrape URL",
            command=self.parent.scrape_single_url,
        )
        self.scrape_button.pack(side="left", padx=(0, 5))

        self.crawl_button = ttk.Button(
            web_row,
            text="Crawl Site",
            command=self.parent.crawl_website,
        )
        self.crawl_button.pack(side="left")

        self.playwright_crawl_button = ttk.Button(
            web_row,
            text="Playwright Crawl",
            command=self.parent.crawl_site_playwright,
        )
        self.playwright_crawl_button.pack(side="left", padx=(5, 0))

        self.manual_verify_button = ttk.Button(
            web_row,
            text="Manual Verify",
            command=self.parent.launch_manual_verification,
        )
        self.manual_verify_button.pack(side="left", padx=(5, 0))

        self.remote_toggle_var = tk.BooleanVar(value=bool(os.getenv("PLAYWRIGHT_REMOTE_ENDPOINT")))
        self.remote_toggle = ttk.Checkbutton(
            web_row,
            text="Attach to existing browser",
            variable=self.remote_toggle_var,
            command=self.parent._toggle_remote_endpoint,
        )
        self.remote_toggle.pack(side="left", padx=(5, 0))

        self.ai_summary_var = tk.BooleanVar(value=True)
        ai_summary_check = ttk.Checkbutton(
            web_row,
            text="AI summary",
            variable=self.ai_summary_var,
        )
        ai_summary_check.pack(side="left", padx=(5, 0))

        remote_controls_row = ttk.Frame(file_ops_frame)
        remote_controls_row.pack(fill="x", padx=5, pady=(2, 2))

        self.launch_chrome_button = ttk.Button(
            remote_controls_row,
            text="Start Chrome",
            command=self.parent.launch_debug_browser,
        )
        self.launch_chrome_button.pack(side="left")

        self.stop_chrome_button = ttk.Button(
            remote_controls_row,
            text="Stop Chrome",
            command=self.parent.stop_debug_browser,
            state="disabled",
        )
        self.stop_chrome_button.pack(side="left", padx=(5, 0))

        self.test_endpoint_button = ttk.Button(
            remote_controls_row,
            text="Test Endpoint",
            command=self.parent.test_remote_debug_endpoint,
        )
        self.test_endpoint_button.pack(side="left", padx=(5, 0))

        self.remote_status_label = ttk.Label(
            file_ops_frame,
            text="No remote browser configured.",
            foreground="gray",
            anchor="w",
        )
        self.remote_status_label.pack(fill="x", padx=10, pady=(0, 5))

        self._sync_remote_controls()

        # Authentication row: Login credentials and authenticated scraping
        auth_row = ttk.Frame(file_ops_frame)
        auth_row.pack(fill="x", padx=5, pady=(0, 5))

        ttk.Label(auth_row, text="Login:").pack(side="left")

        self.username_entry = ttk.Entry(auth_row, width=15)
        self.username_entry.pack(side="left", padx=(5, 2))
        self.username_entry.insert(0, "Username")
        self.username_entry.config(foreground="gray")
        self.username_entry.bind(
            "<FocusIn>",
            lambda e: self._clear_placeholder(self.username_entry, "Username"),
        )
        self.username_entry.bind(
            "<FocusOut>",
            lambda e: self._add_placeholder(self.username_entry, "Username"),
        )

        self.password_entry = ttk.Entry(auth_row, width=15, show="*")
        self.password_entry.pack(side="left", padx=(2, 5))
        self.password_entry.insert(0, "Password")
        self.password_entry.config(show="", foreground="gray")
        self.password_entry.bind(
            "<FocusIn>", lambda e: self._clear_password_placeholder()
        )
        self.password_entry.bind(
            "<FocusOut>", lambda e: self._add_password_placeholder()
        )

        self.login_scrape_button = ttk.Button(
            auth_row,
            text="Login & Scrape",
            command=self.parent.scrape_with_login,
            state="disabled",  # Initially disabled
        )
        self.login_scrape_button.pack(side="left", padx=(0, 5))

        self.analyze_button = ttk.Button(
            auth_row,
            text="Analyze Login",
            command=self.parent.analyze_login_form,
        )
        self.analyze_button.pack(side="left")

        self.navigate_button = ttk.Button(
            auth_row,
            text="Navigate Site",
            command=self.parent.navigate_authenticated_site,
            state="disabled",  # Initially disabled
        )
        self.navigate_button.pack(side="left", padx=(5, 0))

        self.reset_button = ttk.Button(
            auth_row,
            text="Reset",
            command=self.parent.reset_authentication_state,
        )
        self.reset_button.pack(side="left", padx=(5, 0))

        # Management row: File summary and control buttons
        mgmt_row = ttk.Frame(file_ops_frame)
        mgmt_row.pack(fill="x", padx=5, pady=(0, 5))

        # File summary label expands in remaining space
        self.parsed_file_label = ttk.Label(mgmt_row, text="No file loaded.")
        self.parsed_file_label.pack(side="left", fill="x", expand=True)

        self.manage_files_button = ttk.Button(
            mgmt_row, text="Manage Files", command=self.parent.manage_parsed_files
        )
        self.manage_files_button.pack(side="right", padx=(5, 0))

        self.clear_files_button = ttk.Button(
            mgmt_row, text="Clear Files", command=self.parent.clear_uploaded_files
        )
        self.clear_files_button.pack(side="right", padx=(5, 5))

        # Conversation Management Section
        conversation_frame = ttk.LabelFrame(
            content_area, text="Conversation Management"
        )
        conversation_frame.pack(fill="x", pady=5, anchor="n")
        # Use a responsive grid for conversation controls
        conv_inner = ttk.Frame(conversation_frame)
        conv_inner.pack(fill="x", expand=True)
        conv_inner.columnconfigure(4, weight=1)  # Status label stretches

        self.load_conversation_button = ttk.Button(
            conv_inner,
            text="Load Conversation",
            command=self.parent.load_conversation,
        )
        self.load_conversation_button.grid(row=0, column=0, padx=5, pady=5, sticky="w")

        self.save_conversation_button = ttk.Button(
            conv_inner,
            text="Save Conversation",
            command=self.parent.save_conversation,
        )
        self.save_conversation_button.grid(row=0, column=1, padx=5, pady=5, sticky="w")

        # Quick way to return to general chat (no specific prompt bound)
        self.general_chat_button = ttk.Button(
            conv_inner,
            text="General Chat",
            command=self.parent.switch_to_general_chat,
        )
        self.general_chat_button.grid(row=0, column=2, padx=5, pady=5, sticky="w")

        # Info icon popup for full status text
        def _show_status_popup():
            full = getattr(self.conversation_status_label, "full_text", "")
            win = tk.Toplevel(self.parent)
            win.title("Conversation Status")
            win.transient(self.parent)
            win.resizable(True, True)
            # Position near mouse
            x = self.parent.winfo_pointerx() + 12
            y = self.parent.winfo_pointery() + 12
            win.geometry(f"+{x}+{y}")
            frm = ttk.Frame(win, padding=10)
            frm.pack(fill="both", expand=True)
            lbl = ttk.Label(
                frm, text=full or "", justify="left", anchor="w", wraplength=600
            )
            lbl.pack(fill="both", expand=True)
            ttk.Separator(frm, orient="horizontal").pack(fill="x", pady=(8, 6))
            btn = ttk.Button(frm, text="Close", command=win.destroy)
            btn.pack(anchor="e")

        self.info_button = ttk.Button(
            conv_inner, text="ⓘ", width=3, command=_show_status_popup
        )
        self.info_button.grid(row=0, column=3, padx=(0, 5), pady=5, sticky="w")

        # A stretching status area that ellipsizes text and shows a tooltip for full content
        status_frame = ttk.Frame(conv_inner)
        status_frame.grid(row=0, column=4, padx=5, pady=5, sticky="nsew")
        # Ensure this column grows
        conv_inner.columnconfigure(4, weight=1)
        self.conversation_status_label = ttk.Label(
            status_frame, text="No conversation loaded.", anchor="w", justify="left"
        )
        self.conversation_status_label.pack(fill="x", expand=True)

        # Store the full, non-truncated text
        self.conversation_status_label.full_text = "No conversation loaded."

        def _ellipsize_status_text(_event=None):
            # Ellipsize label text based on available width
            try:
                available = max(10, status_frame.winfo_width() - 10)
                full = getattr(self.conversation_status_label, "full_text", "")
                if not full:
                    self.conversation_status_label.config(text="")
                    # Hide icon if no text
                    if hasattr(self, "info_button"):
                        self.info_button.grid_remove()
                    return
                fnt = tkfont.Font(font=self.conversation_status_label["font"])  # type: ignore
                # If full text fits, show it and hide the icon
                if fnt.measure(full) <= available:
                    self.conversation_status_label.config(text=full)
                    if hasattr(self, "info_button"):
                        self.info_button.grid_remove()
                    return
                # Otherwise, show icon only (no truncated text)
                self.conversation_status_label.config(text="")
                if hasattr(self, "info_button"):
                    self.info_button.grid()
            except Exception:
                # Fallback to non-wrapped full text on any error
                self.conversation_status_label.config(
                    text=self.conversation_status_label.cget("text")
                )

        # Bind resize to re-ellipsize
        status_frame.bind("<Configure>", _ellipsize_status_text)
        conv_inner.bind("<Configure>", _ellipsize_status_text)

        # Tooltips: label (when visible) and icon (when shown)
        self._status_tooltip = Tooltip(
            self.conversation_status_label,
            lambda: self.conversation_status_label.full_text,
        )
        self._status_icon_tooltip = Tooltip(
            self.info_button, lambda: self.conversation_status_label.full_text
        )

        # Apply once initially so icon/text are mutually exclusive from the start
        self.parent.after(0, _ellipsize_status_text)

        # Expose helpers for external updates
        self._ellipsize_status_text = _ellipsize_status_text

        def set_conversation_status(text: str):
            self.conversation_status_label.full_text = text or ""
            _ellipsize_status_text()

        self.set_conversation_status = set_conversation_status

        # Context controls: preference + status
        self.context_pref_check = ttk.Checkbutton(
            conv_inner,
            text="Always restrict on switch",
            variable=self.parent.always_restrict_var,
            command=self.parent._save_conversation_state,
        )
        self.context_pref_check.grid(row=0, column=5, padx=5, pady=5, sticky="e")

        self.context_status_label = ttk.Label(conv_inner, text="Context: Shared")
        self.context_status_label.grid(row=0, column=6, padx=5, pady=5, sticky="e")

        chat_frame = ttk.LabelFrame(content_area, text="Chat")
        chat_frame.pack(fill="both", expand=True)

        insights_pane = ttk.PanedWindow(chat_frame, orient="vertical")
        insights_pane.pack(fill="both", expand=True)

        self.crawl_insights_frame = ttk.LabelFrame(insights_pane, text="Live Crawl Insights")
        insights_pane.add(self.crawl_insights_frame, weight=1)

        self.crawl_insights_text = scrolledtext.ScrolledText(
            self.crawl_insights_frame,
            wrap="word",
            height=12,
            state="disabled",
            font=("Proxima Nova Alt", 10),
        )
        self.crawl_insights_text.pack(fill="both", expand=True, padx=10, pady=10)
        self.crawl_insights_text.tag_configure(
            "heading", font=("Proxima Nova Alt", 10, "bold")
        )

        action_frame = ttk.Frame(self.crawl_insights_frame)
        action_frame.pack(fill="x", padx=10, pady=(0, 10))

        self.crawl_followup_button = ttk.Button(
            action_frame,
            text="Ask LLM about current page",
            command=self.parent.ask_llm_about_latest_page,
            state="disabled",
        )
        self.crawl_followup_button.pack(side="left")

        self.crawl_stop_button = ttk.Button(
            action_frame,
            text="Stop crawl",
            command=self.parent.stop_crawl,
            state="disabled",
        )
        self.crawl_stop_button.pack(side="left", padx=(10, 0))

        self.crawl_insights_status = ttk.Label(action_frame, text="No active crawl.")
        self.crawl_insights_status.pack(side="left", padx=(10, 0))

        self.chat_history = scrolledtext.ScrolledText(
            insights_pane,
            wrap="word",
            state="disabled",
            font=("Proxima Nova Alt", 10),
        )
        insights_pane.add(self.chat_history, weight=2)

        # Configure tags for chat message alignment
        self.chat_history.tag_configure("user", justify="right")
        self.chat_history.tag_configure("assistant", justify="left")

        input_frame = ttk.Frame(chat_frame)
        input_frame.pack(fill="x", padx=10, pady=(0, 10))
        input_frame.columnconfigure(0, weight=1)

        # Use a multi-line Text widget for input that auto-expands up to 8 lines
        self.user_input = tk.Text(
            input_frame,
            font=("Proxima Nova Alt", 10),
            wrap="word",
            height=1,
            state="normal",
        )
        self.user_input.grid(row=0, column=0, sticky="nsew", padx=(0, 8))

        # Buttons grouped in a bar so they remain visible on resize
        button_bar = ttk.Frame(input_frame)
        button_bar.grid(row=0, column=1, sticky="e")

        stop_button = ttk.Button(
            button_bar, text="Stop", command=self.parent.stop_response
        )
        stop_button.pack(side="left")

        send_button = ttk.Button(
            button_bar, text="Send", command=self.parent.send_message
        )
        send_button.pack(side="left", padx=(6, 0))

        clear_button = ttk.Button(
            button_bar, text="Clear", command=self.parent.clear_chat
        )
        clear_button.pack(side="left", padx=(6, 0))

        # Enter to send, Shift+Enter for newline
        def _on_return(event):
            # Only send if not holding Shift
            if not (event.state & 0x0001):  # Shift mask
                self.parent.send_message()
                return "break"
            return None

        self.user_input.bind("<Return>", _on_return)

        # Auto-resize on input changes up to 8 lines
        def _auto_resize(_event=None):
            try:
                lines = int(self.user_input.index("end-1c").split(".")[0])
            except Exception:
                lines = 1
            lines = max(1, min(8, lines))
            self.user_input.configure(height=lines)

        self.user_input.bind("<KeyRelease>", _auto_resize)

    def _clear_placeholder(self, entry, placeholder_text):
        """Clear placeholder text on focus."""
        if entry.get() == placeholder_text:
            entry.delete(0, tk.END)
            entry.config(foreground="white")

    def _add_placeholder(self, entry, placeholder_text):
        """Add placeholder text if entry is empty."""
        if not entry.get():
            entry.insert(0, placeholder_text)
            entry.config(foreground="gray")

    def _clear_password_placeholder(self):
        """Clear password placeholder and enable password masking."""
        if self.password_entry.get() == "Password":
            self.password_entry.delete(0, tk.END)
            self.password_entry.config(show="*", foreground="white")

    def _add_password_placeholder(self):
        """Add password placeholder if entry is empty."""
        if not self.password_entry.get():
            self.password_entry.insert(0, "Password")
            self.password_entry.config(show="", foreground="gray")

    def _refresh_endpoint_entry_color(self, *_):
        """Adjust remote endpoint entry color based on its current value."""
        value = self.remote_endpoint_var.get().strip()
        fg = "white" if value else "gray"
        self.endpoint_entry.configure(foreground=fg)

    def _sync_remote_controls(self):
        """Update remote debugging widgets to reflect current state."""
        endpoint = self.remote_endpoint_var.get().strip()
        self.remote_toggle_var.set(bool(endpoint))
        self._refresh_endpoint_entry_color()
        has_process = bool(getattr(self.parent, "chrome_process", None))
        self.launch_chrome_button.config(state="disabled" if has_process else "normal")
        self.stop_chrome_button.config(state="normal" if has_process else "disabled")


logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
# Suppress PyTorch MPS pin_memory warnings on macOS
warnings.filterwarnings(
    "ignore",
    message=r".*'pin_memory' argument is set as true but not supported on MPS.*",
    category=UserWarning,
)


class OllamaGUI(tk.Tk):
    """A GUI for interacting with the Dynamic Ollama Assistant."""

    def __init__(self):
        """Initialize the main application window."""
        super().__init__()
        self.title("Dynamic Ollama Assistant")
        self.geometry("1200x800")
        self.minsize(1000, 600)  # Set minimum window size
        self.parsed_document_content = None
        self.parsed_files = []  # List[dict{name, content}] for multi-file support

        try:
            self.data_by_sheet = load_prompt_catalog(CSV_DIR, CSV_GLOB, EXCEL_GLOB)
        except (FileNotFoundError, IOError) as e:
            messagebox.showerror(
                "Failed to Load Prompts",
                f"An error occurred while loading the prompt catalog (CSV/Excel):\n\n{e}",
            )
            sys.exit(1)

        self.selected_prompt_row = None
        self.is_thinking = False
        self.thinking_animation_id = None
        self.conversation_history = []  # Track conversation messages
        self.system_prompt = None  # Store system prompt to avoid regenerating
        self._current_sheet_name = None
        self._current_row_index = None
        # Context carry-over controls
        self.context_restricted = False  # False => Shared (default), True => Restricted
        self.always_restrict_var = tk.BooleanVar(master=self, value=False)
        # Track which line the current assistant "typing" indicator lives on
        self.current_assistant_line_index: str | None = None
        # Event for stopping in-flight streaming responses
        self.stop_event = threading.Event()
        # Remote Chrome debugging state
        self.chrome_process: Optional[subprocess.Popen[str]] | None = None
        self.chrome_profile_dir: Optional[str] = None
        self.chrome_endpoint: Optional[str] = None
        self._remote_endpoint_healthy = False

        # Initialize UI first
        self.ui = UIComponents(self)
        self.populate_treeview()

        # Then load conversation state after UI is ready
        self._load_conversation_state()

        # Force proper layout calculation after UI creation
        self.update_idletasks()

        # Warm up Ollama model in the background to reduce first-response latency
        threading.Thread(target=self._warm_up_model, daemon=True).start()

        # Prompt to save on close and persist state
        self.protocol("WM_DELETE_WINDOW", self._on_close)

        existing_endpoint = self.ui.remote_endpoint_var.get().strip()
        if existing_endpoint:
            self._update_remote_status(f"Attached to {existing_endpoint}", success=True)
            self._remote_endpoint_healthy = True
        else:
            self._update_remote_status("")
        self.ui._sync_remote_controls()

        # Playwright crawl insight buffer
        self.crawl_paused = False
        self.latest_crawl_entry: Optional[Dict[str, Any]] = None
        self.crawl_result_log: List[Dict[str, Any]] = []
        self.stop_crawl_flag = False

    def _toggle_all(self, open_state: bool):
        """Recursively open or close all items in the treeview."""
        for item in self.ui.prompt_tree.get_children():
            self._toggle_children(item, open_state)

    def _toggle_children(self, parent: str, open_state: bool):
        """Recursively open or close all children of a parent node."""
        self.ui.prompt_tree.item(parent, open=open_state)
        for child in self.ui.prompt_tree.get_children(parent):
            self._toggle_children(child, open_state)

    def expand_all(self):
        """Expand every node in the prompt tree."""
        self._toggle_all(True)

    def collapse_all(self):
        """Collapse every node in the prompt tree."""
        self._toggle_all(False)

    def _update_remote_status(self, message: str, *, success: bool = False):
        """Update the remote connection status label with styling."""
        color = "#39aa56" if success else "#d9534f"
        if not message:
            color = "gray"
            message = "No remote browser configured."
        self.ui.remote_status_label.configure(text=message, foreground=color)

    def _set_remote_endpoint(self, endpoint: str | None):
        """Persist the remote endpoint preference and sync UI state."""
        if endpoint:
            os.environ["PLAYWRIGHT_REMOTE_ENDPOINT"] = endpoint
            self.ui.remote_endpoint_var.set(endpoint)
            self.ui.remote_toggle_var.set(True)
            self._update_remote_status(f"Attached to {endpoint}", success=True)
            self._remote_endpoint_healthy = True
        else:
            os.environ.pop("PLAYWRIGHT_REMOTE_ENDPOINT", None)
            self.ui.remote_endpoint_var.set("")
            self.ui.remote_toggle_var.set(False)
            self._update_remote_status("Playwright will launch its own browser.")
            self._remote_endpoint_healthy = False
        self.ui._sync_remote_controls()

    def _resolve_chrome_path(self) -> Optional[str]:
        """Best-effort lookup for a Chrome executable on macOS/Linux."""
        candidates = [
            "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
            "/Applications/Google Chrome Beta.app/Contents/MacOS/Google Chrome Beta",
            "/Applications/Google Chrome Canary.app/Contents/MacOS/Google Chrome Canary",
            shutil.which("google-chrome"),
            shutil.which("chromium"),
        ]

        for candidate in candidates:
            if candidate and Path(candidate).exists():
                return candidate
        return None

    def _build_debugging_command(self, chrome_path: str, port: int, profile_dir: str) -> list[str]:
        """Construct Chrome launch arguments for remote debugging."""
        return [
            chrome_path,
            f"--remote-debugging-port={port}",
            f"--user-data-dir={profile_dir}",
            "--no-first-run",
            "--disable-popup-blocking",
        ]

    def launch_debug_browser(self):
        """Launch a Chrome instance under remote debugging control."""
        if getattr(self, "chrome_process", None):
            messagebox.showinfo(
                "Chrome Already Running",
                "A debug Chrome session is already active. Use Test Endpoint to verify it.",
            )
            return

        chrome_path = self._resolve_chrome_path()
        if not chrome_path:
            messagebox.showerror(
                "Chrome Not Found",
                "Could not locate Google Chrome. Please install Chrome or set the remote endpoint manually.",
            )
            return

        try:
            port = simpledialog.askinteger(
                "Remote Debugging Port",
                "Choose a remote debugging port (default 9222):",
                initialvalue=9222,
                minvalue=1024,
                maxvalue=65535,
            )
        except Exception:
            port = 9222

        if not port:
            return

        profile_dir = tempfile.mkdtemp(prefix="ollama-chrome-profile-")
        command = self._build_debugging_command(chrome_path, port, profile_dir)

        try:
            proc = subprocess.Popen(command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except Exception as exc:
            shutil.rmtree(profile_dir, ignore_errors=True)
            messagebox.showerror("Launch Failed", f"Failed to start Chrome:\n{exc}")
            return

        endpoint = f"http://localhost:{port}"

        self.chrome_process = proc
        self.chrome_profile_dir = profile_dir
        self.chrome_endpoint = endpoint

        self._set_remote_endpoint(endpoint)
        self._remote_endpoint_healthy = False
        self._update_remote_status(
            f"Chrome running on {endpoint}. Waiting for DevTools…"
        )
        self.ui._sync_remote_controls()

    def stop_debug_browser(self):
        """Terminate the locally launched Chrome remote-debug session."""
        proc = getattr(self, "chrome_process", None)
        if not proc:
            messagebox.showinfo(
                "Chrome Not Running",
                "No locally launched Chrome session is currently active.",
            )
            return

        profile_dir = self.chrome_profile_dir
        self.chrome_process = None
        self.chrome_profile_dir = None
        self.chrome_endpoint = None

        try:
            proc.terminate()
            try:
                proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                proc.kill()
        except Exception as exc:
            logging.warning("Failed to stop Chrome cleanly: %s", exc)

        if profile_dir:
            shutil.rmtree(profile_dir, ignore_errors=True)

        self._set_remote_endpoint(None)
        self._update_remote_status("Chrome debug session stopped.")
        self.ui._sync_remote_controls()

    def test_remote_debug_endpoint(self):
        """Ping the configured remote-debug endpoint to confirm availability."""
        endpoint = self.ui.remote_endpoint_var.get().strip() or getattr(
            self, "chrome_endpoint", ""
        )
        if not endpoint:
            messagebox.showwarning(
                "No Endpoint",
                "Enter a remote debugging endpoint or start Chrome first.",
            )
            return

        parsed = urlparse(endpoint)
        if not parsed.scheme:
            endpoint = f"http://{endpoint}"
            parsed = urlparse(endpoint)

        version_url = f"{parsed.scheme}://{parsed.netloc}/json/version"

        try:
            response = requests.get(version_url, timeout=3)
            response.raise_for_status()
            data = response.json()
            user_agent = data.get("User-Agent", "Unknown UA")
            browser = data.get("Browser", "Unknown Browser")

            status_text = f"Connected to {parsed.netloc} ({browser})"
            self._set_remote_endpoint(f"{parsed.scheme}://{parsed.netloc}")
            self._update_remote_status(status_text, success=True)
            messagebox.showinfo(
                "Endpoint Ready",
                f"Remote debugging endpoint is active.\nBrowser: {browser}\nUser-Agent: {user_agent}",
            )
        except requests.RequestException as exc:
            self._update_remote_status("Failed to connect to remote browser.")
            messagebox.showerror(
                "Endpoint Test Failed",
                f"Could not reach {version_url}.\nDetails: {exc}",
            )
        finally:
            self.ui._sync_remote_controls()

    def _reset_crawl_insights(self):
        """Clear the live crawl insights panel and reset state."""
        self.crawl_result_log.clear()
        self.latest_crawl_entry = None
        self.crawl_paused = False
        self.stop_crawl_flag = False

        text_widget = self.ui.crawl_insights_text
        text_widget.config(state="normal")
        text_widget.delete("1.0", "end")
        text_widget.config(state="disabled")

        self.ui.crawl_followup_button.config(state="disabled")
        self.ui.crawl_stop_button.config(state="disabled")
        self.ui.crawl_insights_status.config(text="No active crawl.")

    def _append_crawl_insight(self, heading: str, body: str):
        """Append formatted text to the live crawl insights panel."""
        text_widget = self.ui.crawl_insights_text
        text_widget.config(state="normal")
        text_widget.insert("end", f"{heading}\n", "heading")
        text_widget.insert("end", f"{body}\n\n")
        text_widget.see("end")
        text_widget.config(state="disabled")

    def _handle_crawl_progress(self, entry: Dict[str, Any]) -> bool:
        """Process each crawl result as it streams in from Playwright."""
        self.latest_crawl_entry = entry
        self.crawl_result_log.append(entry)
        if len(self.crawl_result_log) > 100:
            self.crawl_result_log.pop(0)

        title = entry.get("name", "Crawled Page")
        summary_lines: List[str] = []

        if analysis := entry.get("analysis"):
            summary_lines.append(analysis.strip())
        elif content := entry.get("content"):
            snippet = "\n".join(content.splitlines()[:6])
            if len(content) > 600:
                snippet = snippet[:600] + "…"
            summary_lines.append(snippet)

        if followups := entry.get("followups"):
            formatted = "\n".join(f"- {item}" for item in followups)
            summary_lines.append("Suggested follow-ups:\n" + formatted)

        if candidates := entry.get("candidate_links"):
            sample_links = "\n".join(candidates[:5])
            summary_lines.append("Candidate links:\n" + sample_links)

        body = "\n\n".join(summary_lines) if summary_lines else "No additional details."
        self._append_crawl_insight(title, body)

        self.ui.crawl_followup_button.config(state="normal")
        self.ui.crawl_insights_status.config(text=f"Latest page: {entry.get('url', 'unknown')}")

        if self.stop_crawl_flag:
            self.ui.crawl_insights_status.config(text="Crawl stopping after current page…")
            return False

        return True

    def ask_llm_about_latest_page(self):
        """Seed the chat input with a prompt about the most recent crawled page."""
        if not self.latest_crawl_entry:
            messagebox.showinfo("No Page", "Run a crawl first to populate insights.")
            return

        analysis = self.latest_crawl_entry.get("analysis") or ""
        content = self.latest_crawl_entry.get("content", "")
        url = self.latest_crawl_entry.get("url", "")

        prompt_parts = [f"Let's discuss the page at {url}."]
        if analysis:
            prompt_parts.append("Here is the AI summary:\n" + analysis.strip())
        else:
            snippet = "\n".join(content.splitlines()[:6])
            prompt_parts.append("Here is an excerpt:\n" + snippet)

        prompt_parts.append("Please help me verify key facts and suggest the next steps.")

        user_prompt = "\n\n".join(prompt_parts)
        self.ui.user_input.delete("1.0", "end")
        self.ui.user_input.insert("1.0", user_prompt)
        self.ui.user_input.focus_set()

    def stop_crawl(self):
        """Signal the in-progress crawl loop to stop after current page."""
        if not self.crawl_result_log:
            self.ui.crawl_insights_status.config(text="No crawl is currently running.")
            return
        self.stop_crawl_flag = True
        self.ui.crawl_stop_button.config(state="disabled")
        self.ui.crawl_insights_status.config(text="Stop requested. Finishing current page…")

    def _log_crawl_insights_to_conversation(self, entries: List[Dict[str, Any]]):
        """Append crawl summaries to the conversation history and chat log."""
        if not entries:
            return

        summary_lines: List[str] = ["Playwright crawl summary:"]
        for entry in entries[-5:]:
            url = entry.get("url", "Unknown URL")
            title = entry.get("name", "Page")
            analysis = entry.get("analysis")
            content = entry.get("content", "")
            if analysis:
                snippet = analysis.strip().splitlines()
                snippet_text = snippet[0] if snippet else ""
            else:
                snippet_text = " ".join(content.split())[:160]
            summary_lines.append(f"- **{title}** ({url})\n  {snippet_text}")

        followups_collected: List[str] = []
        for entry in entries[-5:]:
            for follow in entry.get("followups", [])[:2]:
                followups_collected.append(follow)
        if followups_collected:
            summary_lines.append("Next-step questions:")
            for item in followups_collected[:5]:
                summary_lines.append(f"  - {item}")

        message = "\n".join(summary_lines)
        self.conversation_history.append({"role": "assistant", "content": message})
        self.update_chat_history(f"🤖 Assistant: {message}\n\n")
        self._save_conversation_state()

    def clear_placeholder(self, _event=None):
        """Clear placeholder text on focus."""
        if self.ui.search_var.get() == "Search prompts...":
            self.ui.search_var.set("")

    def add_placeholder(self, _event=None):
        """Add placeholder text if entry is empty."""
        if not self.ui.search_var.get():
            self.ui.search_var.set("Search prompts...")

    def on_search(self, *_args):
        """Filter the treeview based on the search query."""
        search_query = self.ui.search_var.get().lower()
        if search_query == "search prompts...":
            search_query = ""
        self.populate_treeview(search_query)

    def populate_treeview(self, search_query=""):
        """Populate the treeview with prompts, optionally filtered by a search query."""
        for item in self.ui.prompt_tree.get_children():
            self.ui.prompt_tree.delete(item)

        for sheet_name, df in self.data_by_sheet.items():
            # Create a copy to avoid modifying the original DataFrame
            filtered_df = df.copy()

            if search_query:
                # Ensure search works on string representations of columns
                mask = filtered_df.apply(
                    lambda row: search_query
                    in str(row["Short Description (PAGE NAME)"]).lower()
                    or search_query in str(row["Sub-Category"]).lower(),
                    axis=1,
                )
                filtered_df = filtered_df[mask]

            if filtered_df.empty:
                continue

            # Collapse categories by default; open only when searching
            sheet_node = self.ui.prompt_tree.insert(
                "",
                "end",
                text=str(sheet_name or "Unnamed Category"),
                open=bool(search_query),
            )

            for sub_cat, group in filtered_df.groupby("Sub-Category"):
                sub_cat_text = str(sub_cat or "Unnamed Sub-Category")
                # Collapse sub-categories by default; open only when searching
                sub_cat_node = self.ui.prompt_tree.insert(
                    sheet_node, "end", text=sub_cat_text, open=bool(search_query)
                )
                for index, row in group.iterrows():
                    page_name = str(
                        row.get("Short Description (PAGE NAME)") or "Unnamed Prompt"
                    )
                    item_id = f"{sheet_name}|{index}"
                    self.ui.prompt_tree.insert(
                        sub_cat_node, "end", text=page_name, iid=item_id
                    )

    def on_prompt_select(self, _event=None):
        """Handle the event when a prompt is selected from the treeview."""
        selected_item = self.ui.prompt_tree.selection()
        if not selected_item:
            return

        item_id = selected_item[0]
        if "|" not in item_id:
            return

        # Parse the selected prompt info
        sheet_name, row_index = item_id.split("|")
        row_index = int(row_index)

        # Check if we're actually switching to a different prompt
        if (
            hasattr(self, "_current_sheet_name")
            and hasattr(self, "_current_row_index")
            and self._current_sheet_name == sheet_name
            and self._current_row_index == row_index
        ):
            # Same prompt selected, no need to clear conversation or UI
            return

        # If there is an active conversation, optionally restrict context per preference
        has_chat_text = self.ui.chat_history.get("1.0", "end").strip()
        if self.conversation_history or has_chat_text:
            if bool(self.always_restrict_var.get()):
                self._clear_conversation_state()
                self.context_restricted = True
            else:
                choice = self._confirm_restrict_context()
                if choice is None:
                    return
                self.context_restricted = bool(choice)
                if choice is True:
                    self._clear_conversation_state()
            self._update_context_status_label()
            self._save_conversation_state()

        # Clear previous prompt's details only when actually switching
        for widget in self.ui.placeholder_frame.winfo_children():
            widget.destroy()
        for widget in self.ui.description_frame.winfo_children():
            widget.destroy()
        self.ui.placeholder_entries.clear()
        self.ui.selected_prompt_label.config(
            text="No Prompt Selected", font=("Proxima Nova Alt", 10, "italic")
        )
        self.selected_prompt_row = None

        # Update the selected prompt label
        page_name = self.ui.prompt_tree.item(item_id, "text")
        self.ui.selected_prompt_label.config(
            text=f"Active: {page_name}", font=("Proxima Nova Alt", 10, "bold")
        )

        sheet_name, index_str = item_id.split("|")
        index = int(index_str)
        self.selected_prompt_row = self.data_by_sheet[sheet_name].loc[index]
        self._current_sheet_name = sheet_name  # Track current sheet for comparison
        self._current_row_index = index  # Track current row index for comparison

        # --- Populate Placeholders and Description ---
        self._populate_details()

    def _clear_prompt_ui(self):
        """Clear any prompt-specific UI (placeholders/description) and selection label."""
        for widget in self.ui.placeholder_frame.winfo_children():
            widget.destroy()
        for widget in self.ui.description_frame.winfo_children():
            widget.destroy()
        self.ui.placeholder_entries.clear()
        self.ui.selected_prompt_label.config(
            text="No Prompt Selected", font=("Proxima Nova Alt", 10, "italic")
        )

    def switch_to_general_chat(self):
        """Switch out of a specific prompt and back to general chat mode."""
        # Ask whether to restrict context if an active conversation exists
        has_chat_text = self.ui.chat_history.get("1.0", "end").strip()
        if self.conversation_history or has_chat_text:
            if bool(self.always_restrict_var.get()):
                self._clear_conversation_state()
                self.context_restricted = True
            else:
                choice = self._confirm_restrict_context()
                if choice is None:  # Cancel switch
                    return
                self.context_restricted = bool(choice)
                if choice is True:
                    self._clear_conversation_state()
            self._update_context_status_label()
            self._save_conversation_state()

        # By default keep current conversation; just detach from any selected prompt
        self.selected_prompt_row = None
        self._current_sheet_name = None
        self._current_row_index = None
        # Clear prompt-specific UI and mark the header
        self._clear_prompt_ui()
        self.ui.selected_prompt_label.config(
            text="General Chat", font=("Proxima Nova Alt", 10, "bold")
        )
        # Force system prompt to be rebuilt on next send (will include parsed doc if any)
        self.system_prompt = None

    def _update_context_status_label(self):
        """Update the context status label to reflect shared/restricted mode."""
        if hasattr(self.ui, "context_status_label"):
            mode = "Restricted" if self.context_restricted else "Shared"
            self.ui.context_status_label.config(text=f"Context: {mode}")

    def _clear_conversation_state(self):
        """Clears the conversation history and resets the chat UI."""
        self.conversation_history.clear()
        self.system_prompt = None
        self.ui.chat_history.config(state="normal")
        self.ui.chat_history.delete("1.0", "end")
        self.ui.chat_history.config(state="disabled")
        # Update status via ellipsis-aware setter
        if hasattr(self.ui, "set_conversation_status"):
            self.ui.set_conversation_status("No conversation loaded.")
        else:
            self.ui.conversation_status_label.config(text="No conversation loaded.")
        self._save_conversation_state()  # Persist the cleared state

    def _save_conversation_state(self):
        """Saves the current conversation history to a state file."""
        state = {
            "conversation_history": self.conversation_history,
            "system_prompt": self.system_prompt,
            "parsed_files": self.parsed_files,
            "parsed_document_content": self.parsed_document_content,
            "context_restricted": self.context_restricted,
            "always_restrict_on_switch": bool(self.always_restrict_var.get()),
        }
        try:
            with open("conversation_state.json", "w", encoding="utf-8") as f:
                json.dump(state, f, indent=2)
        except IOError as e:
            logging.warning("Could not save conversation state: %s", e)

    def _load_conversation_state(self):
        """Loads conversation history from a state file if it exists."""
        try:
            if os.path.exists("conversation_state.json"):
                with open("conversation_state.json", "r", encoding="utf-8") as f:
                    state = json.load(f)
                    self.conversation_history = state.get("conversation_history", [])
                    self.system_prompt = state.get("system_prompt")
                    self.parsed_files = state.get("parsed_files", []) or []
                    self.parsed_document_content = state.get("parsed_document_content")
                    self.context_restricted = bool(
                        state.get("context_restricted", False)
                    )
                    # Load preference; default False
                    self.always_restrict_var.set(
                        bool(state.get("always_restrict_on_switch", False))
                    )

                # Repopulate chat history UI
                self.ui.chat_history.config(state="normal")
                self.ui.chat_history.delete("1.0", "end")
                for msg in self.conversation_history:
                    role = msg.get("role")
                    content = msg.get("content", "")
                    if role == "user":
                        self.update_chat_history(f"👤 User: {content}\n\n")
                    elif role == "assistant":
                        self.update_chat_history(f"🤖 Assistant: {content}\n\n")
                self.ui.chat_history.config(state="disabled")

                # Update parsed files label to reflect loaded state
                self._update_parsed_label_from_state()
                self._update_context_status_label()
        except (IOError, json.JSONDecodeError) as e:
            logging.warning("Could not load conversation state: %s", e)
            self.conversation_history = []

    def _populate_details(self):
        """Populate the UI with details from the selected prompt row."""
        if self.selected_prompt_row is None:
            return

        # --- Description Only --- (align with CSV columns in dynamic_ollama_assistant)
        description = (
            self.selected_prompt_row.get("Description ")
            or self.selected_prompt_row.get("Description")
            or "No description available."
        )

        # Only show the description; do NOT include the mega-prompt or other long fields
        desc_text = str(description).strip()

        desc_label = ttk.Label(
            self.ui.description_frame,
            text=desc_text,
            wraplength=400,
            justify="left",
        )
        desc_label.pack(fill="x")

        # --- Placeholders --- (collect from 'Mega-Prompt' and related fields)
        mega_prompt = self.selected_prompt_row.get("Mega-Prompt", "")
        prompt_name = self.selected_prompt_row.get("Prompt Name", "")
        fields_to_scan = [
            mega_prompt,
            # Restrict placeholder scanning to mega_prompt (and prompt_name for safety)
            prompt_name,
        ]
        placeholders = sorted(
            {ph for field in fields_to_scan for ph in find_placeholders(str(field))}
        )

        # Dynamically create input fields for each placeholder
        for placeholder in placeholders:
            row_frame = ttk.Frame(self.ui.placeholder_frame)
            row_frame.pack(fill="x", pady=2)
            label = ttk.Label(row_frame, text=f"{placeholder}:", width=20)
            label.pack(side="left")
            entry = ttk.Entry(row_frame)
            entry.pack(side="left", fill="x", expand=True)
            self.ui.placeholder_entries[placeholder] = entry

        # Add a button to generate the system prompt
        generate_button = ttk.Button(
            self.ui.placeholder_frame,
            text="Generate System Prompt",
            command=self._generate_system_prompt,
        )
        generate_button.pack(pady=10)

    def _get_combined_document_content(self):
        """Combine all parsed document content from files and scraped content."""
        content_parts = []

        # Add legacy parsed_document_content if it exists
        if getattr(self, "parsed_document_content", None):
            content_parts.append(self.parsed_document_content)

        # Add content from parsed_files (includes scraped content)
        if hasattr(self, "parsed_files") and self.parsed_files:
            for file_data in self.parsed_files:
                if isinstance(file_data, dict) and "content" in file_data:
                    name = file_data.get("name", "Unknown")
                    content = file_data["content"]
                    content_parts.append(f"=== {name} ===\n{content}")

        return "\n\n".join(content_parts) if content_parts else None

    def _generate_system_prompt(self):
        """Build the system prompt using CSV schema and include parsed document if any."""
        if self.selected_prompt_row is None:
            return

        # Collect values from the entry fields
        fill_values = {
            key: entry.get().strip()
            for key, entry in self.ui.placeholder_entries.items()
            if entry.get().strip()
        }
        # Include parsed document content for use by build_system_prompt
        document_content = self._get_combined_document_content()
        if document_content:
            fill_values["parsed_document"] = document_content

        try:
            system_prompt, _ = build_system_prompt(
                self.selected_prompt_row, fill_values
            )
            # Ensure uploaded document is present even if the template omitted it
            if getattr(self, "parsed_document_content", None):
                system_prompt = self._ensure_doc_appended(
                    system_prompt, self.parsed_document_content
                )
            self.system_prompt = system_prompt
            messagebox.showinfo(
                "System Prompt Set",
                "The system prompt has been generated and is ready for your next message.",
            )
        except Exception as e:
            logging.exception("Failed to generate system prompt")
            messagebox.showerror(
                "System Prompt Error",
                f"Failed to generate system prompt: {e}",
            )

    def _load_and_parse_conversation_file(self, file_path):
        """Load chat text and return it along with the expected sidecar path for attachments."""
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        chat_content_start = content.find("👤 User:")
        chat_text = (
            content if chat_content_start == -1 else content[chat_content_start:]
        )

        base, _ = os.path.splitext(file_path)
        sidecar_path = base + ".json"
        return chat_text, sidecar_path

    def _display_loaded_conversation(self, chat_content, file_path, sidecar_path):
        """Update UI with loaded conversation and attach any saved files from sidecar."""
        # Reset current state to tie attachments to this conversation
        self._clear_conversation_state()

        # Load chat text
        self.ui.chat_history.config(state="normal")
        self.ui.chat_history.insert("1.0", chat_content)
        self.ui.chat_history.config(state="disabled")
        if hasattr(self.ui, "set_conversation_status"):
            self.ui.set_conversation_status(f"Loaded: {os.path.basename(file_path)}")
        else:
            self.ui.conversation_status_label.config(
                text=f"Loaded: {os.path.basename(file_path)}"
            )

        # Load attachments from sidecar if present
        try:
            if os.path.exists(sidecar_path):
                with open(sidecar_path, "r", encoding="utf-8") as jf:
                    sc = json.load(jf)
                self.parsed_files = sc.get("parsed_files", []) or []
                self.parsed_document_content = sc.get("parsed_document_content")
                # Invalidate system prompt so next turn reflects attachments
                self.system_prompt = None
                self._update_parsed_label_from_state()
        except (IOError, json.JSONDecodeError) as e:
            logging.warning("Failed to load attachments for conversation: %s", e)

        # Persist
        self._save_conversation_state()
        messagebox.showinfo("Success", "Conversation loaded successfully.")

    def load_conversation(self):
        """Load a conversation from a file."""
        conversations_dir = os.path.join(os.path.dirname(__file__), "conversations")
        file_path = filedialog.askopenfilename(
            title="Load Conversation",
            initialdir=conversations_dir,
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
        )
        if not file_path:
            return

        try:
            chat_content, sidecar_path = self._load_and_parse_conversation_file(
                file_path
            )
            self._display_loaded_conversation(chat_content, file_path, sidecar_path)
        except Exception as e:
            messagebox.showerror("Load Error", f"Failed to load conversation:\n{e}")

    def _confirm_restrict_context(self) -> bool | None:
        """Ask whether to restrict context to the new chat.

        Returns:
            True  -> Restrict context (clear current conversation)
            False -> Keep context across chats (default)
            None  -> Cancel switch
        """
        result = messagebox.askyesnocancel(
            "Restrict Context?",
            "Do you want to restrict context to this chat?\n\n"
            "Yes: Start a fresh context for the new chat.\n"
            "No: Carry over the existing conversation context.",
            icon="question",
        )
        return result

    def clear_chat(self):
        """Clears the current chat conversation."""
        self._clear_conversation_state()

    def save_conversation(self):
        """Save the current conversation to the conversations directory."""

        # Get conversation content from chat history
        conversation_content = self.ui.chat_history.get("1.0", "end").strip()
        if not conversation_content:
            messagebox.showwarning("No Content", "No conversation to save.")
            return

        # Create conversations directory if it doesn't exist
        conversations_dir = os.path.join(os.path.dirname(__file__), "conversations")
        os.makedirs(conversations_dir, exist_ok=True)

        # Ask user for conversation name
        conversation_name = simpledialog.askstring(
            "Save Conversation",
            "Enter a name for this conversation:",
            initialvalue=f"conversation_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}",
        )

        if not conversation_name:
            return

        # Ensure .txt extension
        if not conversation_name.endswith(".txt"):
            conversation_name += ".txt"

        file_path = os.path.join(conversations_dir, conversation_name)

        try:
            # Save human-readable transcript
            with open(file_path, "w", encoding="utf-8") as f:
                f.write("# Conversation Export\n")
                f.write(
                    f"Exported: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
                )
                f.write(conversation_content)

            # Save sidecar with attachments so they can be restored with the transcript
            sidecar_path = os.path.splitext(file_path)[0] + ".json"
            sidecar = {
                "parsed_files": self.parsed_files or [],
                "parsed_document_content": self.parsed_document_content,
            }
            with open(sidecar_path, "w", encoding="utf-8") as jf:
                json.dump(sidecar, jf, indent=2)

            messagebox.showinfo(
                "Saved",
                f"Conversation saved as '{conversation_name}' (with attachments) in conversations folder",
            )
        except IOError as e:
            logging.error("Error saving conversation: %s", e)
            messagebox.showerror("Save Error", f"Failed to save conversation:\n{e}")

    def send_message(self, _event=None):
        """Send a message to the Ollama model and display the response."""
        # Text widget: fetch full text and strip trailing newline
        user_msg = self.ui.user_input.get("1.0", "end-1c").strip()
        if not user_msg or self.is_thinking:
            return

        # Reset any previous stop request before starting a new stream
        self.stop_event.clear()

        # Clear input box (Text indices)
        self.ui.user_input.delete("1.0", "end")
        self.update_chat_history(f"👤 User: {user_msg}\n\n")
        self.conversation_history.append({"role": "user", "content": user_msg})

        # Build an effective system prompt (base prompt + prompt-catalog hints)
        effective_system_prompt = self._build_effective_system_prompt(user_msg)

        self.is_thinking = True
        self.update_chat_history("🤖 Assistant: ")
        # Cancel any stray previous animation just in case
        if getattr(self, "thinking_animation_id", None):
            with contextlib.suppress(Exception):
                self.after_cancel(self.thinking_animation_id)
        self.thinking_animation_id = self.after(100, self._thinking_animation)

        # Use streaming helper that clears the thinking ellipsis on first chunk
        threading.Thread(
            target=self._clear_and_stream_response,
            args=(effective_system_prompt, user_msg),
            daemon=True,
        ).start()

    def _stream_and_process_response(self, user_msg):
        """Helper to stream response from Ollama and update conversation state."""
        try:
            self._stream_ollama_response(user_msg)
        except requests.exceptions.RequestException as e:
            logging.error("Error in Ollama query: %s", e)
            self.after(0, messagebox.showerror, "Error", f"Failed to get response: {e}")
        finally:
            self.is_thinking = False

    def _stream_ollama_response(self, user_msg):
        system_prompt = self.system_prompt or ""
        full_response = ""

        # Stream the response from Ollama
        gen = query_ollama_chat_for_gui(
            model=OLLAMA_MODEL,
            system_prompt=system_prompt,
            user_msg=user_msg,
            conversation_history=self.conversation_history,
        )
        try:
            for chunk in gen:
                if self.stop_event.is_set():
                    # Close underlying HTTP stream via generator close
                    with contextlib.suppress(Exception):
                        gen.close()
                    break
                full_response += chunk
                self.update_chat_history(chunk)
        finally:
            # Ensure generator is closed if we exit early
            with contextlib.suppress(Exception):
                gen.close()

        # Add assistant's response to conversation history
        self.conversation_history.append(
            {"role": "assistant", "content": full_response}
        )
        self.ui.chat_history.config(state=tk.DISABLED)

        # Save conversation state after each assistant response
        self._save_conversation_state()

    def update_chat_history(self, message):
        """Update the chat history in a thread-safe manner."""
        # ... (rest of the method remains the same)
        # Use `after` to schedule the update on the main thread
        self.after(0, self._insert_text, message)

    def _insert_text(self, message):
        """Insert text into the chat history widget with appropriate alignment, tracking the current assistant line."""
        self.ui.chat_history.config(state="normal")

        # Determine the correct tag based on the message sender
        tag = "user" if message.startswith("👤 User:") else "assistant"
        self.ui.chat_history.insert("end", message, tag)
        # If we just inserted the assistant label, remember its line index for the ellipsis animation
        if message.startswith("🤖 Assistant:"):
            # The cursor is at end-1c after insert; store just the line component
            end_index = self.ui.chat_history.index("end-1c")
            self.current_assistant_line_index = end_index.split(".")[0]

        self.ui.chat_history.config(state="disabled")
        self.ui.chat_history.yview("end")

    # ---- Prompt catalog helpers -------------------------------------------------
    def _search_prompts(
        self, query: str, limit: int = 20
    ) -> list[tuple[str, str, str]]:
        """Return top matches from the CSV catalog for a free-text query.

        Each result is a tuple: (sheet_name, sub_category, page_name)
        """
        q = (query or "").strip().lower()
        if not q:
            return []
        tokens = [t for t in q.replace("/", " ").replace("-", " ").split() if t]
        if not tokens:
            return []

        results: list[tuple[int, str, str, str]] = []
        for sheet_name, df in self.data_by_sheet.items():
            for _, row in df.iterrows():
                page = str(row.get("Short Description (PAGE NAME)") or "")
                subc = str(row.get("Sub-Category") or "")
                desc = str(row.get("Description") or row.get("Description ") or "")
                hay = " ".join([page, subc, desc]).lower()
                score = sum(hay.count(tok) for tok in tokens)
                if score > 0:
                    results.append((score, str(sheet_name), subc, page))

        results.sort(key=lambda x: x[0], reverse=True)
        return [(s, c, p) for _, s, c, p in results[:limit]]

    def _build_effective_system_prompt(self, user_msg: str) -> str:
        """Compose the system prompt for this turn, optionally appending catalog hints."""
        # Ensure base system prompt is available (may include parsed document, etc.)
        self._ensure_system_prompt()
        base = self.system_prompt or "You are a helpful assistant."

        matches = self._search_prompts(user_msg, limit=15)
        if not matches:
            return base

        lines = ["\n\nPROMPT CATALOG HINTS (top matches for user's query):"]
        for sheet, subc, page in matches:
            lines.append(f"- [{sheet} > {subc}] {page}")
        return base + "\n" + "\n".join(lines)

    def upload_and_parse_file(self):
        """Handle file upload in the main thread and start parsing in a background thread."""
        file_path = filedialog.askopenfilename(
            title="Select a file to parse",
            filetypes=[
                (
                    "All supported files",
                    "*.pdf *.docx *.html *.md *.pptx *.xlsx *.xls *.csv *.json",
                ),
                ("PDF files", "*.pdf"),
                ("Word documents", "*.docx"),
                ("HTML files", "*.html"),
                ("Markdown files", "*.md"),
                ("PowerPoint presentations", "*.pptx"),
                ("CSV files", "*.csv"),
                ("JSON files", "*.json"),
                ("Excel workbooks", "*.xlsx *.xls"),
                ("All files", "*.*"),
            ],
        )
        if not file_path:
            return

        base = os.path.basename(file_path)
        self.ui.parsed_file_label.config(
            text=f"Parsing {self._truncate_filename(base, max_len=32)}..."
        )
        self._set_parsed_label_tooltip(base)
        self.update_idletasks()

        thread = threading.Thread(target=self._run_file_parsing, args=(file_path,))
        thread.start()

    def upload_and_parse_files(self):
        """Select and parse multiple files in a background thread."""
        logging.info("User initiated file upload dialog")
        file_paths = filedialog.askopenfilenames(
            title="Select files to parse",
            filetypes=[
                (
                    "All supported files",
                    "*.pdf *.docx *.html *.md *.pptx *.xlsx *.xls *.csv *.json",
                ),
                ("PDF files", "*.pdf"),
                ("Word documents", "*.docx"),
                ("HTML files", "*.html"),
                ("Markdown files", "*.md"),
                ("PowerPoint presentations", "*.pptx"),
                ("CSV files", "*.csv"),
                ("JSON files", "*.json"),
                ("Excel workbooks", "*.xlsx *.xls"),
                ("All files", "*.*"),
            ],
        )
        if not file_paths:
            logging.info("User cancelled file selection")
            return
        
        logging.info(f"User selected {len(file_paths)} files for processing:")
        for i, fp in enumerate(file_paths, 1):
            file_size = os.path.getsize(fp) if os.path.exists(fp) else 0
            logging.info(f"  {i}. {os.path.basename(fp)} ({file_size:,} bytes)")
        
        # Create and show progress window
        self._create_progress_window(len(file_paths))
        
        self.ui.parsed_file_label.config(text=f"Parsing {len(file_paths)} file(s)...")
        self.update_idletasks()
        threading.Thread(
            target=self._run_files_parsing, args=(file_paths,), daemon=True
        ).start()

    def _run_file_parsing(self, file_path):
        """Run the file parsing logic in a separate thread."""
        logging.info("Starting to parse file: %s", file_path)
        try:
            self._parse_file_content(file_path)
        except IOError as e:
            logging.error("An error occurred during file parsing: %s", e, exc_info=True)
            self.after(0, self._show_parsing_error, e)

    def _run_files_parsing(self, file_paths):
        """Parse multiple files and aggregate their contents with performance optimizations."""
        logging.info(f"Starting batch file processing for {len(file_paths)} files")
        logging.info(f"Files to process: {[os.path.basename(fp) for fp in file_paths]}")
        
        # Calculate total file size for progress estimation
        total_size = sum(os.path.getsize(fp) if os.path.exists(fp) else 0 for fp in file_paths)
        logging.info(f"Total data to process: {total_size:,} bytes ({total_size / (1024*1024):.1f} MB)")
        
        # Performance optimization: Sort files by size (smaller first for faster feedback)
        file_info = [(fp, os.path.getsize(fp) if os.path.exists(fp) else 0) for fp in file_paths]
        file_info.sort(key=lambda x: x[1])  # Sort by file size
        sorted_files = [fp for fp, _ in file_info]
        
        logging.info("File processing order (optimized by size):")
        for i, (fp, size) in enumerate(file_info, 1):
            logging.info(f"  {i}. {os.path.basename(fp)} ({size:,} bytes)")
        
        aggregated = []
        parsed_list = []
        total_files = len(sorted_files)
        processed_size = 0
        
        for i, fp in enumerate(sorted_files, 1):
            try:
                file_size = os.path.getsize(fp) if os.path.exists(fp) else 0
                logging.info(f"\n{'='*60}")
                logging.info(f"PROCESSING FILE {i}/{total_files}: {os.path.basename(fp)}")
                logging.info(f"File size: {file_size:,} bytes ({file_size / (1024*1024):.2f} MB)")
                logging.info(f"Progress: {((i-1)/total_files)*100:.1f}% complete")
                logging.info(f"Remaining files: {total_files - i + 1}")
                logging.info(f"{'='*60}")
                
                self._update_progress(i-1, total_files, f"Starting {os.path.basename(fp)}...")
                
                start_time = time.time()
                content = self._parse_single_file_collect(fp)
                end_time = time.time()
                processing_time = end_time - start_time
                
                processed_size += file_size
                
                if content:
                    content_length = len(content)
                    logging.info(f"✓ SUCCESS: Extracted {content_length:,} characters from {os.path.basename(fp)}")
                    logging.info(f"Processing time: {processing_time:.2f} seconds")
                    logging.info(f"Processing speed: {file_size / processing_time / 1024:.1f} KB/s" if processing_time > 0 else "Processing speed: N/A")
                    logging.info(f"Data processed so far: {processed_size:,} bytes ({(processed_size/total_size)*100:.1f}% of total)")
                    
                    parsed_list.append(
                        {"name": os.path.basename(fp), "content": content}
                    )
                    header = f"===== FILE: {os.path.basename(fp)} ====="
                    aggregated.append(f"{header}\n\n{content}")
                    
                    self._update_progress(i, total_files, f"✓ Completed {os.path.basename(fp)} ({content_length:,} chars)")
                else:
                    logging.warning(f"✗ FAILED: No content extracted from {os.path.basename(fp)}")
                    logging.info(f"Processing time: {processing_time:.2f} seconds (failed)")
                    self._update_progress(i, total_files, f"✗ Failed {os.path.basename(fp)}")
                    
            except Exception as e:
                logging.error(f"✗ ERROR processing {os.path.basename(fp)}: {str(e)}")
                logging.warning("Skipping file due to parse error: %s (%s)", fp, e)
                self._update_progress(i, total_files, f"✗ Error in {os.path.basename(fp)}")

        # Update state and label on UI thread by appending
        aggregated_text = "\n\n\n".join(aggregated) if aggregated else ""
        total_content_length = len(aggregated_text)
        
        logging.info(f"Batch processing completed: {len(parsed_list)}/{total_files} files processed successfully")
        logging.info(f"Total content extracted: {total_content_length:,} characters")
        
        self.after(0, lambda: self._finalize_file_processing(parsed_list, aggregated_text, len(parsed_list), total_files))

    def _parse_file_content(self, file_path):
        logging.info(f"Starting Docling conversion for: {os.path.basename(file_path)}")
        converter = DocumentConverter()
        result = None
        try:
            logging.info("Initializing DocumentConverter...")
            result = converter.convert(file_path)
            logging.info("Docling conversion completed successfully")
        except Exception as e:
            logging.warning(
                "Docling conversion failed for %s: %s", os.path.basename(file_path), e
            )

        content_parts = []
        if result:
            logging.info("Processing Docling conversion result...")
            # Newer docling may return a list of documents
            if hasattr(result, "documents") and result.documents:
                logging.info(f"Found {len(result.documents)} documents in result")
                content_parts.extend(
                    doc.export_to_markdown() for doc in result.documents
                )
                logging.info("Exported all documents to markdown")
            # Older/other versions may return a single document
            elif hasattr(result, "document") and result.document:
                logging.info("Found single document in result")
                content_parts.append(result.document.export_to_markdown())
                logging.info("Exported document to markdown")
            else:
                logging.warning("Docling result contains no usable documents")

        # Fallback for PPTX using python-pptx if docling produced no content
        if not content_parts and file_path.lower().endswith(".pptx"):
            logging.info("Attempting PPTX fallback processing with python-pptx...")
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                slides_text = []
                slide_count = len(prs.slides)
                logging.info(f"PPTX contains {slide_count} slides")
                
                for i, slide in enumerate(prs.slides, 1):
                    logging.info(f"Processing slide {i}/{slide_count}")
                    slide_shapes = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    slides_text.extend(slide_shapes)
                    
                if slides_text:
                    content_parts.append("\n\n".join(slides_text))
                    logging.info(f"PPTX fallback successful: extracted text from {slide_count} slides")
                else:
                    logging.warning("PPTX fallback found no text content")
            except Exception as e:
                logging.warning("PPTX fallback failed: %s", e)

        # Fallback for Excel using pandas if docling produced no content
        if not content_parts and (
            file_path.lower().endswith(".xlsx") or file_path.lower().endswith(".xls")
        ):
            logging.info("Attempting Excel fallback processing with pandas...")
            try:
                excel = pd.ExcelFile(file_path)
                sheet_count = len(excel.sheet_names)
                logging.info(f"Excel file contains {sheet_count} sheets: {excel.sheet_names}")
                
                sheet_md_parts = []
                for i, sheet in enumerate(excel.sheet_names, 1):
                    logging.info(f"Processing sheet {i}/{sheet_count}: '{sheet}'")
                    try:
                        df = excel.parse(sheet)
                        rows, cols = df.shape
                        logging.info(f"Sheet '{sheet}' has {rows:,} rows and {cols} columns")
                        
                        # Limit overly large tables for responsiveness
                        preview = df.head(200)
                        md = preview.to_markdown(index=False)
                        sheet_md = f"# Sheet: {sheet}\n\n{md}"
                        sheet_md_parts.append(sheet_md)
                        logging.info(f"Successfully processed sheet '{sheet}' (showing first {len(preview)} rows)")
                    except Exception as se:
                        logging.warning(f"Failed to process sheet '{sheet}': {se}")
                        sheet_md_parts.append(
                            f"# Sheet: {sheet}\n\n(Unable to parse sheet: {se})"
                        )
                        
                if sheet_md_parts:
                    content_parts.append("\n\n---\n\n".join(sheet_md_parts))
                    logging.info(
                        "Excel fallback successful: processed %d sheet(s)",
                        len(sheet_md_parts),
                    )
                else:
                    logging.warning("Excel fallback found no processable sheets")
            except Exception as e:
                logging.warning("Excel fallback failed: %s", e)

        # Fallback for CSV using pandas if docling produced no content
        if not content_parts and file_path.lower().endswith(".csv"):
            logging.info("Attempting CSV fallback processing with pandas...")
            try:
                df = pd.read_csv(file_path, on_bad_lines="skip")
                total_rows, total_cols = df.shape
                logging.info(f"CSV file contains {total_rows:,} rows and {total_cols} columns")
                
                preview = df.head(500)
                md = preview.to_markdown(index=False)
                content_parts.append(f"# CSV Preview (first 500 rows)\n\n{md}")
                logging.info(
                    "CSV fallback successful: %s rows previewed out of %s total", len(preview), total_rows
                )
            except Exception as e:
                logging.warning("CSV fallback failed: %s", e)

        # Fallback for JSON: if array of objects, tabularize; else pretty-print
        if not content_parts and file_path.lower().endswith(".json"):
            logging.info("Attempting JSON fallback processing...")
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    
                if isinstance(data, list) and data and isinstance(data[0], dict):
                    logging.info(f"JSON contains {len(data)} records, converting to table format")
                    df = pd.DataFrame(data)
                    rows, cols = df.shape
                    logging.info(f"JSON table has {rows:,} rows and {cols} columns")
                    
                    preview = df.head(500)
                    md = preview.to_markdown(index=False)
                    content_parts.append(
                        f"# JSON Table Preview (first 500 rows)\n\n{md}"
                    )
                    logging.info(f"JSON table fallback successful: showing {len(preview)} rows")
                else:
                    logging.info("JSON contains non-tabular data, formatting as pretty-printed JSON")
                    pretty = json.dumps(data, indent=2, ensure_ascii=False)[:200000]
                    content_parts.append(f"```json\n{pretty}\n```")
                    logging.info(f"JSON pretty-print fallback successful: {len(pretty):,} characters")
                    
                logging.info(
                    "JSON fallback completed for file: %s", os.path.basename(file_path)
                )
            except Exception as e:
                logging.warning("JSON fallback failed: %s", e)

        if content_parts:
            content_text = "\n\n".join(content_parts)
            content_length = len(content_text)
            logging.info(f"Successfully processed {os.path.basename(file_path)}: {content_length:,} characters extracted")
            
            item = {"name": os.path.basename(file_path), "content": content_text}
            header = f"===== FILE: {item['name']} =====\n\n{content_text}"
            # Append on UI thread to keep state consistent
            self.after(0, lambda: self._append_parsed_items([item], header))
            logging.info(
                "Successfully parsed file '%s' (length=%d) and appended to context.",
                item["name"],
                len(content_text),
            )
        else:
            logging.error(f"All processing methods failed for {os.path.basename(file_path)} - no content extracted")
            logging.warning(
                "Failed to extract content from %s. ConversionResult had no document(s).",
                file_path,
            )
            self.after(
                0,
                lambda: self.ui.parsed_file_label.config(
                    text=(
                        "Could not extract content from: "
                        f"{self._truncate_filename(os.path.basename(file_path), max_len=32)}"
                    )
                ),
            )
            self.after(
                0, lambda: self._set_parsed_label_tooltip(os.path.basename(file_path))
            )

    def _create_progress_window(self, total_files):
        """Create and show a progress window for file processing."""
        self.progress_window = tk.Toplevel(self)
        self.progress_window.title("Processing Files")
        self.progress_window.geometry("400x150")
        self.progress_window.resizable(False, False)
        
        # Center the window
        self.progress_window.transient(self)
        self.progress_window.grab_set()
        
        # Progress label
        self.progress_label = ttk.Label(
            self.progress_window, 
            text=f"Processing 0/{total_files} files...",
            font=("Arial", 10)
        )
        self.progress_label.pack(pady=10)
        
        # Progress bar
        self.progress_bar = ttk.Progressbar(
            self.progress_window,
            length=350,
            mode='determinate',
            maximum=total_files
        )
        self.progress_bar.pack(pady=10)
        
        # Status label
        self.progress_status = ttk.Label(
            self.progress_window,
            text="Initializing...",
            font=("Arial", 9),
            foreground="gray"
        )
        self.progress_status.pack(pady=5)
        
        # Cancel button (optional)
        self.progress_cancel = ttk.Button(
            self.progress_window,
            text="Cancel",
            command=self._cancel_processing
        )
        self.progress_cancel.pack(pady=5)
        
        self.processing_cancelled = False
        
    def _update_progress(self, current, total, status_text):
        """Update the progress window with current status."""
        if hasattr(self, 'progress_window') and self.progress_window.winfo_exists():
            self.after(0, lambda: self._update_progress_ui(current, total, status_text))
            
    def _update_progress_ui(self, current, total, status_text):
        """Update progress UI elements on the main thread."""
        if hasattr(self, 'progress_window') and self.progress_window.winfo_exists():
            self.progress_label.config(text=f"Processing {current}/{total} files...")
            self.progress_bar['value'] = current
            self.progress_status.config(text=status_text)
            self.progress_window.update_idletasks()
            
    def _cancel_processing(self):
        """Cancel the file processing operation."""
        self.processing_cancelled = True
        logging.info("User cancelled file processing")
        if hasattr(self, 'progress_window'):
            self.progress_window.destroy()
            
    def _finalize_file_processing(self, parsed_list, aggregated_text, success_count, total_count):
        """Finalize the file processing and update UI."""
        # Close progress window
        if hasattr(self, 'progress_window') and self.progress_window.winfo_exists():
            self.progress_window.destroy()
            
        # Update the parsed items
        self._append_parsed_items(parsed_list, aggregated_text)
        
        # Show completion message
        if success_count == total_count:
            logging.info(f"All {total_count} files processed successfully")
            messagebox.showinfo(
                "Processing Complete", 
                f"Successfully processed all {total_count} files."
            )
        elif success_count > 0:
            logging.info(f"Partial success: {success_count}/{total_count} files processed")
            messagebox.showwarning(
                "Processing Partially Complete", 
                f"Successfully processed {success_count} out of {total_count} files.\n\nCheck the logs for details on failed files."
            )
        else:
            logging.error("No files were successfully processed")
            messagebox.showerror(
                "Processing Failed", 
                "Failed to process any of the selected files.\n\nCheck the logs for error details."
            )

    def _append_parsed_items(self, new_items, aggregated_text: str):
        """Append parsed items and aggregated text to existing state and update label."""
        # Initialize containers if None
        if self.parsed_files is None:
            self.parsed_files = []
        if self.parsed_document_content is None:
            self.parsed_document_content = ""

        # Append items list
        if new_items:
            self.parsed_files.extend(new_items)

        # Append aggregated text with proper separator
        if aggregated_text:
            if self.parsed_document_content:
                self.parsed_document_content += "\n\n\n" + aggregated_text
            else:
                self.parsed_document_content = aggregated_text

        # Invalidate any previously cached system prompt so the next turn includes the new document
        # contents via `_ensure_system_prompt()` or explicit generation.
        self.system_prompt = None

        # Update label and persist
        self._update_parsed_label_from_state()
        self._save_conversation_state()

    def _update_parsed_label_from_state(self):
        """Update the parsed_file_label based on current parsed_files/document content."""
        total = len(self.parsed_files or [])
        if total > 0:
            # Truncate each displayed filename to avoid pushing UI buttons out of view
            names = ", ".join(
                self._truncate_filename(item["name"], max_len=24)
                for item in self.parsed_files[:3]
            )
            more = "" if total <= 3 else f" (+{total-3} more)"
            self.ui.parsed_file_label.config(
                text=f"Loaded: {total} file(s): {names}{more}"
            )
            # Tooltip with full filenames
            tooltip_text = "\n".join(
                item.get("name", "<unnamed>") for item in (self.parsed_files or [])
            )
            self._set_parsed_label_tooltip(tooltip_text)
        elif self.parsed_document_content:
            self.ui.parsed_file_label.config(text="Loaded: cached document")
            self._set_parsed_label_tooltip("Cached document content loaded")
        else:
            self.ui.parsed_file_label.config(text="No file loaded.")
            self._set_parsed_label_tooltip("")

    def manage_parsed_files(self):
        """Open a simple dialog to remove specific uploaded files."""
        if not (self.parsed_files or []):
            messagebox.showinfo("Manage Files", "No files to manage.")
            return

        top = tk.Toplevel(self)
        top.title("Manage Files")
        top.geometry("450x300")
        top.transient(self)
        top.grab_set()

        ttk.Label(top, text="Select files to remove:").pack(
            anchor="w", padx=10, pady=(10, 5)
        )

        frame = ttk.Frame(top)
        frame.pack(fill="both", expand=True, padx=10, pady=5)

        scrollbar = ttk.Scrollbar(frame, orient="vertical")
        listbox = tk.Listbox(frame, selectmode="extended", yscrollcommand=scrollbar.set)
        scrollbar.config(command=listbox.yview)
        listbox.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        for item in self.parsed_files:
            listbox.insert("end", item.get("name", "<unnamed>"))

        btn_frame = ttk.Frame(top)
        btn_frame.pack(fill="x", padx=10, pady=10)

        def remove_selected():
            sel = listbox.curselection()
            if not sel:
                return
            names_to_remove = {listbox.get(i) for i in sel}
            self.parsed_files = [
                it
                for it in (self.parsed_files or [])
                if it.get("name") not in names_to_remove
            ]
            # Update UI/state
            self._update_parsed_label_from_state()
            self._save_conversation_state()
            # Refresh list
            listbox.delete(0, "end")
            for it in self.parsed_files:
                listbox.insert("end", it.get("name", "<unnamed>"))
            if not self.parsed_files:
                top.destroy()

        remove_btn = ttk.Button(
            btn_frame, text="Remove Selected", command=remove_selected
        )
        remove_btn.pack(side="left")

        close_btn = ttk.Button(btn_frame, text="Close", command=top.destroy)
        close_btn.pack(side="right")

    def _truncate_filename(self, filename: str, max_len: int = 24) -> str:
        """Truncate a filename with middle ellipsis, preserving extension when present."""
        if not filename or len(filename) <= max_len:
            return filename
        root, ext = os.path.splitext(filename)
        # Leave room for ext and ellipsis
        budget = max(5, max_len - len(ext) - 3)
        return (
            f"{self._truncate_middle(root, budget)}…{ext}"
            if ext
            else self._truncate_middle(filename, max_len)
        )

    def _truncate_middle(self, text: str, max_len: int) -> str:
        """Truncate the middle of a string with an ellipsis to fit max_len."""
        if len(text) <= max_len:
            return text
        half = (max_len - 1) // 2
        return f"{text[:half]}…{text[-(max_len - 1 - half):]}"

    # --- Tooltip helpers for showing full filenames on hover ---
    def _set_parsed_label_tooltip(self, text: str):
        """Attach or update a tooltip for the parsed files label."""
        # Create lazily
        if not hasattr(self, "_parsed_files_tooltip"):
            self._parsed_files_tooltip = ToolTip(self.ui.parsed_file_label, text=text)
        else:
            self._parsed_files_tooltip.text = text or ""

    def _parse_single_file_collect(self, file_path) -> str:
        """Parse a single file and return its extracted text with performance optimizations."""
        file_ext = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        
        logging.info(f"  → Starting processing pipeline for {file_ext} file: {os.path.basename(file_path)}")
        
        # Performance optimization: Skip very large files or use streaming for them
        if file_size > 100 * 1024 * 1024:  # 100MB threshold
            logging.warning(f"  → Large file detected ({file_size / (1024*1024):.1f} MB). Using optimized processing...")
            return self._parse_large_file_optimized(file_path)
        
        # Performance optimization: Use faster fallback for known simple formats first
        if file_ext in ['.txt', '.md', '.csv', '.json']:
            logging.info(f"  → Fast-track processing for simple format: {file_ext}")
            fast_result = self._try_fast_processing(file_path, file_ext)
            if fast_result:
                return fast_result
        
        converter = DocumentConverter()
        result = None
        try:
            logging.info(f"  → Step 1/4: Initializing Docling DocumentConverter for {file_ext} format")
            start_time = time.time()
            
            # Performance optimization: Configure Docling for faster processing with timeout
            import signal
            
            def timeout_handler(signum, frame):
                raise TimeoutError("Docling conversion timed out")
            
            # Set timeout for Docling conversion (5 minutes max)
            signal.signal(signal.SIGALRM, timeout_handler)
            signal.alarm(300)  # 5 minutes timeout
            
            try:
                result = converter.convert(file_path)
                signal.alarm(0)  # Cancel timeout
                
                conversion_time = time.time() - start_time
                logging.info(f"  → Step 1/4: ✓ Docling conversion completed in {conversion_time:.2f}s")
            except TimeoutError:
                signal.alarm(0)  # Cancel timeout
                raise TimeoutError("Docling conversion exceeded 5 minute timeout")
                
        except Exception as e:
            logging.warning(
                "  → Step 1/4: ✗ Docling conversion failed for %s: %s", os.path.basename(file_path), e
            )
            logging.info(f"  → Will attempt fallback processing methods...")

        content_parts = []
        if result:
            logging.info(f"  → Step 2/4: Processing Docling conversion result...")
            if hasattr(result, "documents") and result.documents:
                doc_count = len(result.documents)
                logging.info(f"  → Found {doc_count} document(s) in Docling result")
                for i, doc in enumerate(result.documents, 1):
                    logging.info(f"  → Exporting document {i}/{doc_count} to markdown...")
                    markdown_content = doc.export_to_markdown()
                    content_parts.append(markdown_content)
                    logging.info(f"  → Document {i}/{doc_count}: {len(markdown_content):,} characters exported")
                logging.info(f"  → Step 2/4: ✓ All {doc_count} documents exported to markdown")
            elif hasattr(result, "document") and result.document:
                logging.info(f"  → Found single document in Docling result")
                markdown_content = result.document.export_to_markdown()
                content_parts.append(markdown_content)
                logging.info(f"  → Step 2/4: ✓ Single document exported: {len(markdown_content):,} characters")
            else:
                logging.warning(f"  → Step 2/4: ✗ Docling result contains no usable documents")
        else:
            logging.info(f"  → Step 2/4: No Docling result to process, moving to fallback methods")

        if not content_parts and file_path.lower().endswith(".pptx"):
            logging.info(f"  → Step 3/4: Attempting PPTX fallback processing with python-pptx...")
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                slides_text = []
                slide_count = len(prs.slides)
                logging.info(f"  → PPTX file contains {slide_count} slides")
                
                for i, slide in enumerate(prs.slides, 1):
                    logging.info(f"  → Processing slide {i}/{slide_count}...")
                    slide_shapes = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    if slide_shapes:
                        slides_text.extend(slide_shapes)
                        logging.info(f"  → Slide {i}: extracted {len(slide_shapes)} text elements")
                    else:
                        logging.info(f"  → Slide {i}: no text content found")
                        
                if slides_text:
                    combined_text = "\n\n".join(slides_text)
                    content_parts.append(combined_text)
                    logging.info(f"  → Step 3/4: ✓ PPTX fallback successful: {len(combined_text):,} characters from {slide_count} slides")
                else:
                    logging.warning(f"  → Step 3/4: ✗ PPTX fallback found no text content in any slides")
            except Exception as e:
                logging.warning("  → Step 3/4: ✗ PPTX fallback failed for %s: %s", os.path.basename(file_path), e)

        # Excel fallback with performance optimizations
        if not content_parts and (
            file_path.lower().endswith(".xlsx") or file_path.lower().endswith(".xls")
        ):
            logging.info(f"  → Step 3/4: Attempting Excel fallback processing with pandas...")
            try:
                excel = pd.ExcelFile(file_path)
                sheet_count = len(excel.sheet_names)
                logging.info(f"  → Excel file contains {sheet_count} sheets: {excel.sheet_names}")
                
                # Performance optimization: Limit number of sheets processed for large files
                file_size_mb = file_size / (1024 * 1024)
                max_sheets = 10 if file_size_mb > 20 else sheet_count
                if sheet_count > max_sheets:
                    logging.info(f"  → Large Excel file ({file_size_mb:.1f} MB), processing first {max_sheets} sheets only")
                    sheets_to_process = excel.sheet_names[:max_sheets]
                else:
                    sheets_to_process = excel.sheet_names
                
                sheet_md_parts = []
                for i, sheet in enumerate(sheets_to_process, 1):
                    logging.info(f"  → Processing sheet {i}/{len(sheets_to_process)}: '{sheet}'...")
                    try:
                        # Performance optimization: Read only first N rows for large files
                        nrows = 1000 if file_size_mb > 10 else None
                        df = excel.parse(sheet, nrows=nrows)
                        rows, cols = df.shape
                        logging.info(f"  → Sheet '{sheet}': {rows:,} rows × {cols} columns")
                        
                        # Performance optimization: Adaptive preview size
                        preview_size = min(100 if file_size_mb > 20 else 200, rows)
                        preview = df.head(preview_size)
                        preview_rows = len(preview)
                        logging.info(f"  → Creating markdown preview for first {preview_rows} rows...")
                        
                        md = preview.to_markdown(index=False)
                        sheet_md = f"# Sheet: {sheet}\n\n{md}"
                        sheet_md_parts.append(sheet_md)
                        logging.info(f"  → Sheet '{sheet}': ✓ {len(sheet_md):,} characters generated")
                    except Exception as se:
                        logging.warning(f"  → Sheet '{sheet}': ✗ Failed to process: {se}")
                        sheet_md_parts.append(
                            f"# Sheet: {sheet}\n\n(Unable to parse sheet: {se})"
                        )
                        
                if sheet_md_parts:
                    combined_content = "\n\n---\n\n".join(sheet_md_parts)
                    content_parts.append(combined_content)
                    logging.info(f"  → Step 3/4: ✓ Excel fallback successful: {len(combined_content):,} characters from {len(sheet_md_parts)} sheets")
                else:
                    logging.warning(f"  → Step 3/4: ✗ Excel fallback found no processable sheets")
            except Exception as e:
                logging.warning("  → Step 3/4: ✗ Excel fallback failed for %s: %s", os.path.basename(file_path), e)

        # CSV fallback with performance optimizations
        if not content_parts and file_path.lower().endswith(".csv"):
            logging.info(f"  → Step 3/4: Attempting CSV fallback processing with pandas...")
            try:
                logging.info(f"  → Reading CSV file...")
                
                # Performance optimization: Use chunking for large CSV files
                file_size_mb = file_size / (1024 * 1024)
                if file_size_mb > 50:  # For files larger than 50MB
                    logging.info(f"  → Large CSV detected ({file_size_mb:.1f} MB), using chunked reading...")
                    chunk_size = 1000
                    df_chunks = pd.read_csv(file_path, chunksize=chunk_size, on_bad_lines="skip")
                    first_chunk = next(df_chunks)
                    total_cols = len(first_chunk.columns)
                    preview = first_chunk.head(500)
                    logging.info(f"  → CSV chunk processed: {len(preview)} rows × {total_cols} columns (chunked mode)")
                else:
                    df = pd.read_csv(file_path, on_bad_lines="skip")
                    total_rows, total_cols = df.shape
                    logging.info(f"  → CSV contains {total_rows:,} rows × {total_cols} columns")
                    
                    # Performance optimization: Limit preview size based on file size
                    preview_size = min(500 if file_size_mb < 10 else 200, total_rows)
                    logging.info(f"  → Creating preview of first {preview_size} rows...")
                    preview = df.head(preview_size)
                
                logging.info(f"  → Converting to markdown format...")
                md = preview.to_markdown(index=False)
                csv_content = f"# CSV Preview (first {len(preview)} rows)\n\n{md}"
                content_parts.append(csv_content)
                
                logging.info(f"  → Step 3/4: ✓ CSV fallback successful: {len(csv_content):,} characters ({len(preview)} rows processed)")
            except Exception as e:
                logging.warning("  → Step 3/4: ✗ CSV fallback failed for %s: %s", os.path.basename(file_path), e)

        # JSON fallback
        if not content_parts and file_path.lower().endswith(".json"):
            logging.info(f"  → Step 3/4: Attempting JSON fallback processing...")
            try:
                logging.info(f"  → Reading JSON file...")
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    
                logging.info(f"  → Analyzing JSON structure...")
                if isinstance(data, list) and data and isinstance(data[0], dict):
                    record_count = len(data)
                    logging.info(f"  → JSON contains {record_count:,} records (tabular format detected)")
                    
                    logging.info(f"  → Converting to DataFrame...")
                    df = pd.DataFrame(data)
                    rows, cols = df.shape
                    logging.info(f"  → DataFrame: {rows:,} rows × {cols} columns")
                    
                    preview_size = min(500, rows)
                    logging.info(f"  → Creating table preview of first {preview_size} rows...")
                    preview = df.head(preview_size)
                    
                    md = preview.to_markdown(index=False)
                    json_content = f"# JSON Table Preview (first {preview_size} rows)\n\n{md}"
                    content_parts.append(json_content)
                    logging.info(f"  → Step 3/4: ✓ JSON table fallback successful: {len(json_content):,} characters ({preview_size}/{record_count} records)")
                else:
                    logging.info(f"  → JSON contains non-tabular data, formatting as pretty-printed JSON")
                    pretty = json.dumps(data, indent=2, ensure_ascii=False)[:200000]
                    json_content = f"```json\n{pretty}\n```"
                    content_parts.append(json_content)
                    logging.info(f"  → Step 3/4: ✓ JSON pretty-print fallback successful: {len(json_content):,} characters")
                    
            except Exception as e:
                logging.warning("  → Step 3/4: ✗ JSON fallback failed for %s: %s", os.path.basename(file_path), e)
        
        # PDF fallback using PyPDF2 or pdfplumber when Docling fails
        if not content_parts and file_path.lower().endswith(".pdf"):
            logging.info(f"  → Step 3/4: Attempting PDF fallback processing...")
            
            # Try PyPDF2 first (faster but less accurate)
            try:
                logging.info(f"  → Trying PyPDF2 for basic text extraction...")
                import PyPDF2
                
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    num_pages = len(pdf_reader.pages)
                    logging.info(f"  → PDF contains {num_pages} pages")
                    
                    text_parts = []
                    max_pages = min(50, num_pages)  # Limit to first 50 pages for performance
                    
                    for i in range(max_pages):
                        try:
                            page = pdf_reader.pages[i]
                            text = page.extract_text()
                            if text.strip():
                                text_parts.append(f"# Page {i+1}\n\n{text}")
                                logging.info(f"  → Page {i+1}/{max_pages}: {len(text)} characters extracted")
                            else:
                                logging.info(f"  → Page {i+1}/{max_pages}: no text found")
                        except Exception as pe:
                            logging.warning(f"  → Page {i+1} extraction failed: {pe}")
                    
                    if text_parts:
                        pdf_content = "\n\n---\n\n".join(text_parts)
                        if max_pages < num_pages:
                            pdf_content += f"\n\n[... Showing first {max_pages} of {num_pages} pages ...]"
                        content_parts.append(pdf_content)
                        logging.info(f"  → Step 3/4: ✓ PyPDF2 fallback successful: {len(pdf_content):,} characters from {len(text_parts)} pages")
                    else:
                        logging.warning(f"  → PyPDF2 extracted no text from PDF")
                        
            except ImportError:
                logging.info(f"  → PyPDF2 not available, trying pdfplumber...")
                
                # Try pdfplumber as secondary fallback
                try:
                    import pdfplumber
                    
                    with pdfplumber.open(file_path) as pdf:
                        num_pages = len(pdf.pages)
                        logging.info(f"  → PDF contains {num_pages} pages (pdfplumber)")
                        
                        text_parts = []
                        max_pages = min(20, num_pages)  # More conservative limit for pdfplumber
                        
                        for i in range(max_pages):
                            try:
                                page = pdf.pages[i]
                                text = page.extract_text()
                                if text and text.strip():
                                    text_parts.append(f"# Page {i+1}\n\n{text}")
                                    logging.info(f"  → Page {i+1}/{max_pages}: {len(text)} characters extracted")
                                else:
                                    logging.info(f"  → Page {i+1}/{max_pages}: no text found")
                            except Exception as pe:
                                logging.warning(f"  → Page {i+1} extraction failed: {pe}")
                        
                        if text_parts:
                            pdf_content = "\n\n---\n\n".join(text_parts)
                            if max_pages < num_pages:
                                pdf_content += f"\n\n[... Showing first {max_pages} of {num_pages} pages ...]"
                            content_parts.append(pdf_content)
                            logging.info(f"  → Step 3/4: ✓ pdfplumber fallback successful: {len(pdf_content):,} characters from {len(text_parts)} pages")
                        else:
                            logging.warning(f"  → pdfplumber extracted no text from PDF")
                            
                except ImportError:
                    logging.warning(f"  → No PDF fallback libraries available (PyPDF2, pdfplumber)")
                    # Create a basic message for unsupported PDF
                    pdf_content = f"# PDF Processing Failed\n\nThis PDF file could not be processed. Docling conversion failed and no fallback PDF libraries are available.\n\nTo enable PDF fallback processing, install:\n- PyPDF2: `pip install PyPDF2`\n- pdfplumber: `pip install pdfplumber`"
                    content_parts.append(pdf_content)
                    logging.info(f"  → Step 3/4: Added PDF processing failure notice")
                    
            except Exception as e:
                logging.warning(f"  → Step 3/4: ✗ PDF fallback failed: {e}")
                # Still provide a helpful message
                pdf_content = f"# PDF Processing Error\n\nFailed to process PDF file: {str(e)}\n\nThis may be due to:\n- Encrypted/password-protected PDF\n- Corrupted PDF file\n- Complex PDF layout\n- Missing dependencies"
                content_parts.append(pdf_content)
                logging.info(f"  → Step 3/4: Added PDF error message")

        # Final result compilation
        if content_parts:
            final_content = "\n\n".join(content_parts)
            logging.info(f"  → Step 4/4: ✓ Content compilation successful")
            logging.info(f"  → Final result: {len(final_content):,} total characters from {len(content_parts)} content part(s)")
            logging.info(f"  → Processing complete for {os.path.basename(file_path)}")
            return final_content
        else:
            logging.error(f"  → Step 4/4: ✗ No content extracted from any processing method")
            logging.error(f"  → Processing failed for {os.path.basename(file_path)}")
            return ""
    
    def _try_fast_processing(self, file_path, file_ext):
        """Try fast processing for simple file formats."""
        try:
            if file_ext in ['.txt', '.md']:
                logging.info(f"  → Fast text processing for {file_ext}...")
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                if content.strip():
                    logging.info(f"  → Fast processing successful: {len(content):,} characters")
                    return content
                    
            elif file_ext == '.csv':
                logging.info(f"  → Fast CSV processing...")
                # Quick CSV processing for smaller files
                file_size = os.path.getsize(file_path)
                if file_size < 10 * 1024 * 1024:  # Less than 10MB
                    df = pd.read_csv(file_path, nrows=1000, on_bad_lines="skip")
                    rows, cols = df.shape
                    md = df.head(min(200, rows)).to_markdown(index=False)
                    content = f"# CSV Quick Preview ({rows} rows × {cols} columns)\n\n{md}"
                    logging.info(f"  → Fast CSV processing successful: {len(content):,} characters")
                    return content
                    
            elif file_ext == '.json':
                logging.info(f"  → Fast JSON processing...")
                file_size = os.path.getsize(file_path)
                if file_size < 5 * 1024 * 1024:  # Less than 5MB
                    with open(file_path, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                    
                    if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                        df = pd.DataFrame(data[:500])  # Limit to first 500 records
                        md = df.head(100).to_markdown(index=False)
                        content = f"# JSON Quick Preview (first 100 records)\n\n{md}"
                    else:
                        content = f"```json\n{json.dumps(data, indent=2)[:50000]}\n```"
                    
                    logging.info(f"  → Fast JSON processing successful: {len(content):,} characters")
                    return content
                    
        except Exception as e:
            logging.info(f"  → Fast processing failed for {file_ext}: {e}, falling back to standard processing")
            
        return None
    
    def _parse_large_file_optimized(self, file_path):
        """Optimized processing for very large files (>100MB)."""
        file_ext = os.path.splitext(file_path)[1].lower()
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        
        logging.info(f"  → Large file optimization for {file_size_mb:.1f} MB {file_ext} file")
        
        try:
            if file_ext in ['.txt', '.md']:
                # Read first 1MB of text files
                logging.info(f"  → Reading first 1MB of large text file...")
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(1024 * 1024)  # 1MB
                content += "\n\n[... File truncated due to size ...]" if len(content) == 1024 * 1024 else ""
                logging.info(f"  → Large text file processing: {len(content):,} characters extracted")
                return content
                
            elif file_ext == '.csv':
                # Process first 5000 rows of large CSV
                logging.info(f"  → Processing first 5000 rows of large CSV...")
                df = pd.read_csv(file_path, nrows=5000, on_bad_lines="skip")
                rows, cols = df.shape
                preview = df.head(200)
                md = preview.to_markdown(index=False)
                content = f"# Large CSV Sample ({rows:,} rows × {cols} columns - showing first 200)\n\n{md}\n\n[... File truncated due to size ...]"
                logging.info(f"  → Large CSV processing: {len(content):,} characters extracted")
                return content
                
            elif file_ext in ['.xlsx', '.xls']:
                # Process only first sheet, limited rows
                logging.info(f"  → Processing first sheet of large Excel file...")
                excel = pd.ExcelFile(file_path)
                first_sheet = excel.sheet_names[0]
                df = excel.parse(first_sheet, nrows=1000)
                rows, cols = df.shape
                preview = df.head(100)
                md = preview.to_markdown(index=False)
                content = f"# Large Excel Sample - Sheet: {first_sheet} ({rows:,} rows × {cols} columns - showing first 100)\n\n{md}\n\n[... File truncated due to size ...]"
                logging.info(f"  → Large Excel processing: {len(content):,} characters extracted")
                return content
                
            elif file_ext == '.pdf':
                # For large PDFs, try basic text extraction with limits
                logging.info(f"  → Processing large PDF with basic extraction...")
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        num_pages = len(pdf_reader.pages)
                        
                        # For large PDFs, only process first 10 pages
                        max_pages = min(10, num_pages)
                        text_parts = []
                        
                        for i in range(max_pages):
                            try:
                                page = pdf_reader.pages[i]
                                text = page.extract_text()
                                if text.strip():
                                    # Limit text per page for large files
                                    text = text[:5000] + "..." if len(text) > 5000 else text
                                    text_parts.append(f"# Page {i+1}\n\n{text}")
                            except Exception:
                                continue
                        
                        if text_parts:
                            content = "\n\n---\n\n".join(text_parts)
                            content += f"\n\n[... Large PDF: Showing first {max_pages} of {num_pages} pages ...]"
                            logging.info(f"  → Large PDF processing: {len(content):,} characters from {len(text_parts)} pages")
                            return content
                        
                except ImportError:
                    pass
                
                # Fallback message for large PDFs
                content = f"# Large PDF Notice\n\nThis PDF file ({file_size_mb:.1f} MB) is too large for full processing.\nFor better results with large PDFs, consider:\n- Using a dedicated PDF processing tool\n- Splitting the PDF into smaller files\n- Installing PyPDF2 or pdfplumber for basic text extraction"
                logging.info(f"  → Large PDF: provided notice message")
                return content
                
            else:
                # For other large files, return a message
                content = f"# Large File Notice\n\nThis {file_ext} file ({file_size_mb:.1f} MB) is too large for full processing.\nPlease consider splitting it into smaller files or using a more specific processing tool."
                logging.info(f"  → Large file skipped: {file_ext} format not supported for large files")
                return content
                
        except Exception as e:
            logging.error(f"  → Large file optimization failed: {e}")
            return f"# Large File Processing Error\n\nFailed to process large {file_ext} file ({file_size_mb:.1f} MB): {str(e)}"

    def clear_uploaded_files(self):
        """Clear any uploaded/parsed files and their aggregated content."""
        previous_count = len(self.parsed_files) if self.parsed_files else 0
        self.parsed_files = []
        self.parsed_document_content = None
        self.ui.parsed_file_label.config(text="No file loaded.")
        self._set_parsed_label_tooltip("")
        self._save_conversation_state()
        logging.info(f"Cleared {previous_count} uploaded files")

    def _show_parsing_error(self, error):
        """Display parsing error in a thread-safe way."""
        self.ui.parsed_file_label.config(text="Failed to parse file.")
        messagebox.showerror(
            "Parsing Error", f"An error occurred while parsing the file:\n{error}"
        )

    def _clear_and_stream_response(self, system_prompt, user_msg):
        """Keep thinking until first chunk, then clear and stream the response from Ollama."""
        try:
            full_response = ""
            first_chunk = True
            start_time = time.time()
            gen = query_ollama_chat_for_gui(
                model=OLLAMA_MODEL,
                system_prompt=system_prompt,
                user_msg=user_msg,
                conversation_history=self.conversation_history,
            )
            for chunk in gen:
                if self.stop_event.is_set():
                    with contextlib.suppress(Exception):
                        gen.close()
                    break
                if first_chunk:
                    # Stop the animation and clear the ellipsis once the first chunk arrives
                    self.is_thinking = False

                    def _clear_ellipsis():
                        # Clear only the dots on the current assistant line, if tracked
                        if not self.current_assistant_line_index:
                            return
                        self.ui.chat_history.config(state="normal")
                        line_index = self.current_assistant_line_index
                        start_pos = f"{line_index}.0 + 13c"
                        line_end_pos = f"{line_index}.end"
                        current_line_content = self.ui.chat_history.get(
                            start_pos, line_end_pos
                        )
                        if current_line_content.strip() in {".", "..", "..."}:
                            self.ui.chat_history.delete(start_pos, line_end_pos)
                        self.ui.chat_history.config(state="disabled")

                    self.after(0, _clear_ellipsis)
                    ttfb = time.time() - start_time
                    logging.info("TTFB (time-to-first-byte): %.2fs", ttfb)
                    first_chunk = False

                full_response += chunk
                self.update_chat_history(chunk)

            total_time = time.time() - start_time
            logging.info("Total response time: %.2fs", total_time)
            # Append the full response to the conversation history
            self.conversation_history.append(
                {"role": "assistant", "content": full_response}
            )
            # Save conversation state after each assistant response
            self._save_conversation_state()
            # Reset typing line tracking now that the response is finished
            self.current_assistant_line_index = None
            # Add spacing after assistant reply for clearer UI separation
            self.update_chat_history("\n\n")
        except requests.exceptions.RequestException as e:
            logging.error("Error during streaming: %s", e)
            self.after(0, messagebox.showerror, "Error", f"Failed to get response: {e}")
        finally:
            # Ensure thinking state is reset even if an error occurs
            self.is_thinking = False
            # Re-enable input
            with contextlib.suppress(tk.TclError):
                self.ui.user_input.config(state="normal")
                self.ui.user_input.focus_set()

    def stop_response(self):
        """Signal any in-flight streaming to stop and tidy up UI state."""
        self.stop_event.set()
        # Stop thinking animation immediately
        with contextlib.suppress(Exception):
            if getattr(self, "thinking_animation_id", None):
                self.after_cancel(self.thinking_animation_id)
                self.thinking_animation_id = None
        self.is_thinking = False

    def _warm_up_model(self):
        """Send a tiny background request to keep the model warm."""
        try:
            # Minimal system prompt and user ping
            sys_prompt = "You are a helpful assistant. Reply with a single dot."
            for _ in query_ollama_chat_for_gui(
                model=OLLAMA_MODEL,
                system_prompt=sys_prompt,
                user_msg="ping",
                conversation_history=[],
            ):
                # Stop after first small chunk
                break
            logging.info("Model warm-up completed.")
        except requests.exceptions.RequestException as e:
            logging.debug("Model warm-up skipped or failed: %s", e)

    def _thinking_animation(self, dot_count=0):
        """Animate the thinking ellipsis in the chat window."""
        if not self.is_thinking:
            return

        ellipsis = "." * ((dot_count % 3) + 1)
        # Only touch the line for the current assistant typing indicator
        if self.current_assistant_line_index:
            line_index = self.current_assistant_line_index
            start_pos = f"{line_index}.0 + 13c"
            line_end_pos = f"{line_index}.end"
            self.ui.chat_history.config(state="normal")
            self.ui.chat_history.delete(start_pos, line_end_pos)
            self.ui.chat_history.insert(start_pos, ellipsis, "assistant")
            self.ui.chat_history.config(state="disabled")

        self.thinking_animation_id = self.after(
            500, lambda: self._thinking_animation(dot_count + 1)
        )

    def _ensure_system_prompt(self):
        """Ensure there is a system prompt ready, auto-including parsed document if available."""
        if self.system_prompt:
            return

        fill_values = {}
        if document_content := self._get_combined_document_content():
            fill_values["parsed_document"] = document_content

        try:
            if self.selected_prompt_row is not None:
                # Build using the selected row schema so placeholders are respected
                system_prompt, _ = build_system_prompt(
                    self.selected_prompt_row, fill_values
                )
                # Guarantee the uploaded document is appended even if template missed it
                if fill_values.get("parsed_document"):
                    system_prompt = self._ensure_doc_appended(
                        system_prompt, fill_values["parsed_document"]
                    )
                self.system_prompt = system_prompt
            elif fill_values.get("parsed_document"):
                # Fallback minimal system prompt that includes the document
                doc = fill_values["parsed_document"]
                self.system_prompt = (
                    "You are a helpful assistant. Use the provided document to answer.\n\n"
                    "USER-PROVIDED DOCUMENT:\n---BEGIN DOCUMENT---\n"
                    f"{doc}\n"
                    "---END DOCUMENT---\n"
                )
            else:
                # Final fallback minimal prompt
                self.system_prompt = "You are a helpful assistant."
        except Exception as e:
            logging.warning("Auto system prompt generation failed: %s", e)
            # Ensure we still have something
            if not self.system_prompt:
                self.system_prompt = "You are a helpful assistant."

    def _ensure_doc_appended(self, prompt: str, document: str) -> str:
        """Append the uploaded document block to the prompt if it's not already present."""
        marker_tokens = [
            "---BEGIN DOCUMENT---",
            "USER-PROVIDED DOCUMENT",
            "BEGIN DOCUMENT",
        ]
        if any(tok in (prompt or "") for tok in marker_tokens):
            return prompt
        # Don't duplicate if a very long substring already present
        if document and document[:200] in (prompt or ""):
            return prompt
        return (
            (prompt or "You are a helpful assistant.")
            + "\n\nUSER-PROVIDED DOCUMENT:\n---BEGIN DOCUMENT---\n"
            + (document or "")
            + "\n---END DOCUMENT---\n"
        )

    def scrape_single_url(self):
        """Scrape content from a single URL and add it to parsed files."""
        url = self.ui.url_entry.get().strip()

        if not url or url == self.ui.url_placeholder:
            messagebox.showerror("Error", "Please enter a URL to scrape.")
            return

        if not validate_url(url):
            messagebox.showerror("Error", "Please enter a valid URL.")
            return

        try:
            # Show progress
            self.ui.scrape_button.config(text="Scraping...", state="disabled")
            self.update_idletasks()

            # Scrape the content
            scraped_data = scrape_web_content(url)

            # Add to parsed files
            self.parsed_files.append(scraped_data)

            # Update UI
            self._update_parsed_label_from_state()
            self._save_conversation_state()

            # Clear URL entry
            self.ui.url_entry.delete(0, tk.END)

            messagebox.showinfo(
                "Success", f"Successfully scraped: {scraped_data['name']}"
            )

        except Exception as e:
            messagebox.showerror("Scraping Error", f"Failed to scrape URL:\n{str(e)}")
        finally:
            self.ui.scrape_button.config(text="Scrape URL", state="normal")

    def crawl_website(self):
        """Crawl multiple pages from a website and add them to parsed files."""
        url = self.ui.url_entry.get().strip()

        if not url or url == self.ui.url_placeholder:
            messagebox.showerror("Error", "Please enter a URL to crawl.")
            return

        if not validate_url(url):
            messagebox.showerror("Error", "Please enter a valid URL.")
            return

        # Ask user for crawl settings
        max_pages = simpledialog.askinteger(
            "Crawl Settings",
            "Maximum number of pages to crawl:",
            initialvalue=5,
            minvalue=1,
            maxvalue=20,
        )

        if not max_pages:
            return

        try:
            # Show progress
            self.ui.crawl_button.config(text="Crawling...", state="disabled")
            self.update_idletasks()

            # Crawl the website
            scraped_pages = crawl_website(url, max_pages=max_pages)

            if not scraped_pages:
                messagebox.showwarning(
                    "No Content", "No pages could be scraped from the website."
                )
                return

            # Add all scraped pages to parsed files
            self.parsed_files.extend(scraped_pages)

            # Update UI
            self._update_parsed_label_from_state()
            self._save_conversation_state()

            # Clear URL entry
            self.ui.url_entry.delete(0, tk.END)

            messagebox.showinfo(
                "Success",
                f"Successfully crawled {len(scraped_pages)} pages from the website.",
            )

        except Exception as e:
            messagebox.showerror(
                "Crawling Error", f"Failed to crawl website:\n{str(e)}"
            )
        finally:
            self.ui.crawl_button.config(text="Crawl Site", state="normal")

    def crawl_site_playwright(self):
        """Crawl pages using Playwright to handle dynamic or protected sites."""
        url = self.ui.url_entry.get().strip()

        if not url or url == self.ui.url_placeholder:
            messagebox.showerror("Error", "Please enter a URL to crawl.")
            return

        max_pages = simpledialog.askinteger(
            "Crawl Settings",
            "Maximum number of pages to crawl (Playwright):",
            initialvalue=5,
            minvalue=1,
            maxvalue=20,
        )

        if not max_pages:
            return

        try:
            self._reset_crawl_insights()
            self.ui.crawl_insights_status.config(text="Playwright crawl in progress…")
            self.ui.crawl_stop_button.config(state="normal")
            self.ui.playwright_crawl_button.config(text="Crawling...", state="disabled")
            self.update_idletasks()

            username = self.ui.username_entry.get().strip()
            password = self.ui.password_entry.get().strip()
            creds_provided = (
                username
                and password
                and username != "Username"
                and password != "Password"
            )

            login_selectors = (
                self.ui.login_selectors
                if getattr(self.ui, "login_analyzed", False)
                else None
            )

            include_ai_summary = bool(self.ui.ai_summary_var.get())

            def _progress(entry: Dict[str, Any]) -> bool:
                return self._handle_crawl_progress(entry)

            results = playwright_crawl_sync(
                url,
                max_pages=max_pages,
                same_domain_only=True,
                username=username if creds_provided else None,
                password=password if creds_provided else None,
                login_selectors=login_selectors,
                include_ai_summary=include_ai_summary,
                progress_callback=_progress,
            )

            success = [res for res in results if "Error" not in res.get("name", "")]

            if not success:
                messagebox.showwarning(
                    "No Content", "No pages were successfully crawled with Playwright."
                )
                self.ui.crawl_insights_status.config(text="Crawl finished without content.")
                return

            self.parsed_files.extend(success)
            self._update_parsed_label_from_state()
            self._save_conversation_state()

            if include_ai_summary:
                self._log_crawl_insights_to_conversation(success)

            messagebox.showinfo(
                "Success",
                f"Crawled {len(success)} page(s) using Playwright."
                + (" Summaries generated." if include_ai_summary else ""),
            )

            self.ui.crawl_insights_status.config(
                text=f"Crawl complete. {len(success)} page(s) captured."
            )

        except Exception as exc:
            messagebox.showerror("Crawl Error", f"Playwright crawl failed:\n{str(exc)}")
            self.ui.crawl_insights_status.config(text="Crawl failed.")
        finally:
            self.stop_crawl_flag = False
            self.ui.crawl_stop_button.config(state="disabled")
            self.ui.playwright_crawl_button.config(
                text="Playwright Crawl", state="normal"
            )

    def launch_manual_verification(self):
        """Open a dialog instructing the user to complete verification manually."""
        top = tk.Toplevel(self)
        top.title("Manual Verification Required")
        top.geometry("420x260")
        top.transient(self)
        top.grab_set()

        ttk.Label(
            top,
            text=(
                "A bot-detection challenge is blocking automation.\n"
                "If you enabled 'Attach to existing browser', solve the challenge "
                "in that Chrome window. Otherwise, you can open a new verification "
                "browser below.\nOnce the site is accessible, click 'Session ready'."
            ),
            wraplength=380,
            justify="left",
        ).pack(padx=15, pady=(20, 15), anchor="w")

        ttk.Button(
            top,
            text="Open Browser",
            command=lambda: self._open_manual_browser(top),
        ).pack(pady=(0, 10))

        ttk.Button(
            top,
            text="Session ready",
            command=lambda: self._on_manual_ready(top),
        ).pack(pady=(0, 5))

        ttk.Button(top, text="Cancel", command=top.destroy).pack(pady=(0, 10))

    def _open_manual_browser(self, dialog):
        try:
            from playwright.sync_api import sync_playwright

            url = self.ui.url_entry.get().strip()
            if not url or url == self.ui.url_placeholder:
                messagebox.showerror("Error", "Please enter a URL to open.")
                return

            def _launch():
                with sync_playwright() as p:
                    browser = p.chromium.launch(headless=False)
                    page = browser.new_page()
                    page.goto(url)
                    messagebox.showinfo(
                        "Browser Launched",
                        "A Playwright-controlled browser has been opened.\n"
                        "Complete the verification and close the browser when finished.",
                    )
                    browser.close()

            threading.Thread(target=_launch, daemon=True).start()
        except Exception as exc:
            messagebox.showerror(
                "Launch Error",
                f"Failed to open manual verification browser:\n{exc}",
            )

    def _on_manual_ready(self, dialog):
        dialog.destroy()
        messagebox.showinfo(
            "Session Ready",
            "Great! The session cookies should now allow automated crawling again."
        )

    def _toggle_remote_endpoint(self):
        attach = bool(self.ui.remote_toggle_var.get())
        endpoint = self.ui.remote_endpoint_var.get().strip()

        if attach:
            if not endpoint:
                messagebox.showerror(
                    "Missing Endpoint",
                    "Enter the remote debugging URL (e.g. http://localhost:9222) before enabling attachment.",
                )
                self.ui.remote_toggle_var.set(False)
                return
            os.environ["PLAYWRIGHT_REMOTE_ENDPOINT"] = endpoint
            messagebox.showinfo(
                "Remote Attachment Enabled",
                f"Playwright will now connect to the browser at {endpoint}.",
            )
        else:
            os.environ.pop("PLAYWRIGHT_REMOTE_ENDPOINT", None)
            messagebox.showinfo(
                "Remote Attachment Disabled",
                "Playwright will launch its own browser instances again.",
            )

    def scrape_with_login(self):
        """Scrape content from a site requiring authentication."""
        url = self.ui.url_entry.get().strip()
        username = self.ui.username_entry.get().strip()
        password = self.ui.password_entry.get().strip()

        if not url or url == self.ui.url_placeholder:
            messagebox.showerror("Error", "Please enter a URL to scrape.")
            return

        if not validate_url(url):
            messagebox.showerror("Error", "Please enter a valid URL.")
            return

        if not username or username == "Username":
            messagebox.showerror("Error", "Please enter a username.")
            return

        if not password or password == "Password":
            messagebox.showerror("Error", "Please enter a password.")
            return

        try:
            # Show progress
            self.ui.login_scrape_button.config(text="Logging in...", state="disabled")
            self.update_idletasks()

            # Use stored selectors if available, otherwise let the scraper detect them
            login_selectors = (
                self.ui.login_selectors if self.ui.login_analyzed else None
            )

            # Scrape with authentication
            result = scrape_with_login_sync(url, username, password, login_selectors)

            # Check for verification requirements
            if result.get("requires_manual_verification"):
                verification_info = result.get("verification_info", {})
                self._show_verification_required_dialog(verification_info, url)
                return

            if "Error" in result["name"] or "Failed" in result["name"]:
                messagebox.showerror("Authentication Error", result["content"])
                return

            # Mark as successfully authenticated
            self.ui.authenticated_session = True
            self._update_auth_button_states()

            # Add to parsed files
            self.parsed_files.append(result)
            self._update_parsed_label_from_state()
            self._save_conversation_state()

            # Clear credentials for security
            self._clear_credentials()

            messagebox.showinfo(
                "Success",
                f"Successfully scraped authenticated content from {result['name']}",
            )

        except Exception as e:
            messagebox.showerror(
                "Authentication Error", f"Failed to scrape with login:\n{str(e)}"
            )
        finally:
            self.ui.login_scrape_button.config(text="Login & Scrape", state="normal")
            self._update_auth_button_states()

    def analyze_login_form(self):
        """Analyze a page to help identify login form elements."""
        url = self.ui.url_entry.get().strip()

        if not url or url == self.ui.url_placeholder:
            messagebox.showerror("Error", "Please enter a URL to analyze.")
            return

        if not validate_url(url):
            messagebox.showerror("Error", "Please enter a valid URL.")
            return

        try:
            # Show progress
            self.ui.analyze_button.config(text="Analyzing...", state="disabled")
            self.update_idletasks()

            # Analyze the login form
            selectors = analyze_login_form_sync(url)

            if "error" in selectors:
                if selectors.get("manual_mode"):
                    # Show manual mode dialog for 403 errors
                    self._show_manual_selector_dialog(selectors, url)
                else:
                    messagebox.showerror(
                        "Analysis Error", f"Analysis failed:\n{selectors['error']}"
                    )
                    self.ui.login_analyzed = False
                    self.ui.login_selectors = None
            else:
                self._show_login_analysis(selectors, url)
                # Mark as analyzed and enable Login & Scrape button
                self.ui.login_analyzed = True
                self.ui.login_selectors = selectors
                self._update_auth_button_states()

        except Exception as e:
            messagebox.showerror(
                "Analysis Error", f"Failed to analyze login form:\n{str(e)}"
            )
            self.ui.login_analyzed = False
            self.ui.login_selectors = None
        finally:
            self.ui.analyze_button.config(text="Analyze Login", state="normal")
            self._update_auth_button_states()

    def _show_manual_selector_dialog(self, error_info, url):
        """Show manual selector entry dialog for sites that block automated analysis."""
        dialog = tk.Toplevel(self)
        dialog.title("Manual Selector Entry - Site Blocks Analysis")
        dialog.transient(self)
        dialog.grab_set()

        # Center the dialog
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (700 // 2)
        y = (dialog.winfo_screenheight() // 2) - (500 // 2)
        dialog.geometry(f"700x500+{x}+{y}")

        main_frame = ttk.Frame(dialog, padding="10")
        main_frame.pack(fill="both", expand=True)

        # Error explanation
        ttk.Label(
            main_frame,
            text="⚠️ Site Blocks Automated Analysis",
            font=("Arial", 14, "bold"),
            foreground="#ff6b35",
        ).pack(pady=(0, 10))

        # Error details
        error_frame = ttk.LabelFrame(main_frame, text="Error Details", padding="10")
        error_frame.pack(fill="x", pady=(0, 15))

        error_text = tk.Text(
            error_frame, height=3, wrap="word", background="#2b2b2b", foreground="white"
        )
        error_text.pack(fill="x")
        error_text.insert(
            "1.0", f"{error_info['error']}\n\n{error_info.get('suggestion', '')}"
        )
        error_text.config(state="disabled")

        # Manual entry section
        manual_frame = ttk.LabelFrame(
            main_frame, text="Manual CSS Selector Entry", padding="10"
        )
        manual_frame.pack(fill="both", expand=True, pady=(0, 15))

        # Instructions
        ttk.Label(
            manual_frame,
            text="Enter CSS selectors manually by inspecting the login page:",
            font=("Arial", 10, "bold"),
        ).pack(anchor="w", pady=(0, 10))

        # Selector entries
        entries = {}
        common_selectors = error_info.get("common_selectors", {})

        for field, placeholder in [
            ("username", "Username/Email field selector"),
            ("password", "Password field selector"),
            ("submit", "Submit button selector"),
        ]:
            field_frame = ttk.Frame(manual_frame)
            field_frame.pack(fill="x", pady=5)

            ttk.Label(field_frame, text=f"{field.title()}:", width=12).pack(side="left")
            entry = ttk.Entry(field_frame, width=50)
            entry.pack(side="left", fill="x", expand=True, padx=(5, 0))

            # Pre-fill with common selector suggestions
            if field in common_selectors:
                entry.insert(0, common_selectors[field].split(",")[0].strip())

            entries[field] = entry

        # Help text
        help_frame = ttk.LabelFrame(
            manual_frame, text="How to Find Selectors", padding="10"
        )
        help_frame.pack(fill="x", pady=(10, 0))

        help_text = tk.Text(
            help_frame, height=6, wrap="word", background="#2b2b2b", foreground="white"
        )
        help_text.pack(fill="x")
        help_text.insert(
            "1.0",
            """1. Open the login page in your browser
2. Right-click on the username field → "Inspect Element"
3. Copy the CSS selector (right-click element → Copy → Copy selector)
4. Repeat for password field and submit button
5. Common patterns: input[name="email"], #password, button[type="submit"]""",
        )
        help_text.config(state="disabled")

        # Buttons
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill="x", pady=(10, 0))

        def use_manual_selectors():
            selectors = {}
            for field, entry in entries.items():
                value = entry.get().strip()
                if value:
                    selectors[field] = value

            if len(selectors) >= 2:  # At least username and password
                self.ui.login_analyzed = True
                self.ui.login_selectors = selectors
                self._update_auth_button_states()
                dialog.destroy()
                messagebox.showinfo(
                    "Success", "Manual selectors saved! You can now use Login & Scrape."
                )
            else:
                messagebox.showerror(
                    "Error", "Please enter at least username and password selectors."
                )

        def open_browser():
            import webbrowser

            webbrowser.open(url)

        ttk.Button(button_frame, text="Open Login Page", command=open_browser).pack(
            side="left"
        )
        ttk.Button(button_frame, text="Cancel", command=dialog.destroy).pack(
            side="right"
        )
        ttk.Button(
            button_frame, text="Use These Selectors", command=use_manual_selectors
        ).pack(side="right", padx=(0, 10))

    def _show_verification_required_dialog(self, verification_info, url):
        """Show dialog when CAPTCHA or human verification is detected."""
        dialog = tk.Toplevel(self)
        dialog.title("Human Verification Required")
        dialog.transient(self)
        dialog.grab_set()

        # Center the dialog
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (600 // 2)
        y = (dialog.winfo_screenheight() // 2) - (400 // 2)
        dialog.geometry(f"600x400+{x}+{y}")

        main_frame = ttk.Frame(dialog, padding="10")
        main_frame.pack(fill="both", expand=True)

        # Warning header
        ttk.Label(
            main_frame,
            text="🤖 Human Verification Detected",
            font=("Arial", 14, "bold"),
            foreground="#ff6b35",
        ).pack(pady=(0, 10))

        # Verification details
        details_frame = ttk.LabelFrame(
            main_frame, text="Verification Details", padding="10"
        )
        details_frame.pack(fill="x", pady=(0, 15))

        details_text = tk.Text(
            details_frame,
            height=4,
            wrap="word",
            background="#2b2b2b",
            foreground="white",
        )
        details_text.pack(fill="x")

        details_content = f"Site: {verification_info.get('current_url', url)}\n"
        details_content += (
            f"Page Title: {verification_info.get('page_title', 'Unknown')}\n\n"
        )

        if verification_info.get("content_matches"):
            details_content += (
                f"Detected: {', '.join(verification_info['content_matches'][:3])}"
            )

        details_text.insert("1.0", details_content)
        details_text.config(state="disabled")

        # Instructions
        instructions_frame = ttk.LabelFrame(
            main_frame, text="What This Means", padding="10"
        )
        instructions_frame.pack(fill="both", expand=True, pady=(0, 15))

        instructions_text = tk.Text(
            instructions_frame,
            height=8,
            wrap="word",
            background="#2b2b2b",
            foreground="white",
        )
        instructions_text.pack(fill="x")
        instructions_text.insert(
            "1.0",
            """This site requires human verification (CAPTCHA, "I'm not a robot", etc.) which cannot be automated.

Options to proceed:
1. Manual Login: Complete the verification manually in your browser, then try scraping
2. Session Import: If you have valid cookies/session data, import them
3. Alternative Approach: Some sites offer API access or different login methods

Common verification types detected:
• reCAPTCHA ("I'm not a robot" checkbox)
• hCaptcha (image/text challenges)  
• Cloudflare bot protection
• Custom verification challenges""",
        )
        instructions_text.config(state="disabled")

        # Buttons
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill="x", pady=(10, 0))

        def open_site():
            import webbrowser

            webbrowser.open(url)

        def try_anyway():
            """Allow user to attempt login despite verification detection."""
            dialog.destroy()
            # Continue with normal login flow
            messagebox.showinfo(
                "Proceeding",
                "Attempting login despite verification detection. Manual intervention may be required.",
            )

        ttk.Button(button_frame, text="Open Site Manually", command=open_site).pack(
            side="left"
        )
        ttk.Button(button_frame, text="Cancel", command=dialog.destroy).pack(
            side="right"
        )
        ttk.Button(button_frame, text="Try Login Anyway", command=try_anyway).pack(
            side="right", padx=(0, 10)
        )

    def _show_login_analysis(self, selectors, url):
        """Show login form analysis results in a dialog."""
        dialog = tk.Toplevel(self)
        dialog.title("Login Form Analysis")
        dialog.transient(self)
        dialog.grab_set()

        # Center the dialog
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (600 // 2)
        y = (dialog.winfo_screenheight() // 2) - (400 // 2)
        dialog.geometry(f"600x400+{x}+{y}")

        main_frame = ttk.Frame(dialog, padding="10")
        main_frame.pack(fill="both", expand=True)

        ttk.Label(
            main_frame,
            text=f"Login Form Analysis for: {url}",
            font=("TkDefaultFont", 10, "bold"),
        ).pack(pady=(0, 10))

        # Results display
        results_frame = ttk.LabelFrame(main_frame, text="Detected Elements")
        results_frame.pack(fill="both", expand=True, pady=(0, 10))

        text_widget = scrolledtext.ScrolledText(results_frame, wrap="word", height=15)
        text_widget.pack(fill="both", expand=True, padx=10, pady=10)

        content = "The AI assistant identified the following login form elements:\n\n"
        for field, selector in selectors.items():
            content += f"{field.title()}: {selector}\n"

        content += "\n" + "=" * 50 + "\n"
        content += "You can use these selectors for authenticated scraping.\n"
        content += (
            "If the selectors don't work, try manually inspecting the page source."
        )

        text_widget.insert("1.0", content)
        text_widget.config(state="disabled")

        # Buttons
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill="x")

        ttk.Button(button_frame, text="Close", command=dialog.destroy).pack(
            side="right"
        )

        def use_selectors():
            messagebox.showinfo(
                "Info",
                "Selectors saved! You can now use 'Login & Scrape' with your credentials.",
            )
            dialog.destroy()

        ttk.Button(
            button_frame, text="Use These Selectors", command=use_selectors
        ).pack(side="right", padx=(0, 10))

    def _update_auth_button_states(self):
        """Update button states based on authentication workflow state."""
        if hasattr(self.ui, "login_analyzed") and self.ui.login_analyzed:
            # Login form has been analyzed - enable Login & Scrape
            self.ui.login_scrape_button.config(state="normal")
            # Disable Analyze Login to prevent re-analysis
            self.ui.analyze_button.config(state="disabled")
        else:
            # No login analysis yet - disable Login & Scrape
            self.ui.login_scrape_button.config(state="disabled")
            # Enable Analyze Login
            self.ui.analyze_button.config(state="normal")

        # If we have an authenticated session, enable navigation
        if hasattr(self.ui, "authenticated_session") and self.ui.authenticated_session:
            self.ui.navigate_button.config(state="normal")
        else:
            self.ui.navigate_button.config(state="disabled")

    def _clear_credentials(self):
        """Clear username and password fields."""
        self.ui.username_entry.delete(0, tk.END)
        self.ui.username_entry.insert(0, "Username")

        self.ui.password_entry.config(show="")
        self.ui.password_entry.delete(0, tk.END)
        self.ui.password_entry.insert(0, "Password")

    def reset_authentication_state(self):
        """Reset all authentication fields and button states."""
        # Clear URL field and restore placeholder
        self.ui.url_entry.delete(0, tk.END)
        self.ui.url_entry.insert(0, self.ui.url_placeholder)
        self.ui.url_entry.config(foreground="gray")

        # Clear credentials
        self._clear_credentials()

        # Reset authentication state flags
        if hasattr(self.ui, "login_analyzed"):
            self.ui.login_analyzed = False
        if hasattr(self.ui, "authenticated_session"):
            self.ui.authenticated_session = False
        if hasattr(self.ui, "login_selectors"):
            self.ui.login_selectors = None

        # Update button states to initial state
        self.ui.analyze_button.config(state="normal")
        self.ui.login_scrape_button.config(state="disabled")
        self.ui.navigate_button.config(state="disabled")

        # Update status
        if hasattr(self.ui, "set_conversation_status"):
            self.ui.set_conversation_status("Authentication state reset")
        else:
            self.ui.conversation_status_label.config(text="Authentication state reset")

        # Clear any stored session data
        try:
            import os

            session_file = "scraper_sessions.json"
            if os.path.exists(session_file):
                os.remove(session_file)
        except Exception:
            pass  # Ignore errors when clearing session file

    def _on_close(self):
        """Handle application close: persist state, cancel animations, and destroy the window."""
        with contextlib.suppress(Exception):
            self._save_conversation_state()

        # Cancel any pending thinking animation
        with contextlib.suppress(Exception):
            if getattr(self, "thinking_animation_id", None):
                self.after_cancel(self.thinking_animation_id)
                self.thinking_animation_id = None

        # Optionally, you could prompt the user to confirm close if desired.
        with contextlib.suppress(Exception):
            self.destroy()

    def navigate_authenticated_site(self):
        """Show dialog for navigating authenticated site and scraping additional pages."""
        if (
            not hasattr(self.ui, "authenticated_session")
            or not self.ui.authenticated_session
        ):
            messagebox.showerror("Error", "Please login first using 'Login & Scrape'.")
            return

        dialog = tk.Toplevel(self)
        dialog.title("Navigate Authenticated Site")
        dialog.geometry("600x400")
        dialog.transient(self)
        dialog.grab_set()

        # Center the dialog
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (600 // 2)
        y = (dialog.winfo_screenheight() // 2) - (400 // 2)
        dialog.geometry(f"600x400+{x}+{y}")

        # Main frame
        main_frame = ttk.Frame(dialog)
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)

        # Instructions
        instructions = ttk.Label(
            main_frame,
            text="Enter URLs to navigate and scrape (one per line):",
            font=("TkDefaultFont", 10, "bold"),
        )
        instructions.pack(anchor="w", pady=(0, 5))

        # URL input area
        url_frame = ttk.Frame(main_frame)
        url_frame.pack(fill="both", expand=True, pady=(0, 10))

        url_text = tk.Text(url_frame, height=8, wrap="word")
        url_scrollbar = ttk.Scrollbar(
            url_frame, orient="vertical", command=url_text.yview
        )
        url_text.configure(yscrollcommand=url_scrollbar.set)

        url_text.pack(side="left", fill="both", expand=True)
        url_scrollbar.pack(side="right", fill="y")

        # Add placeholder text
        placeholder = "https://example.com/page1\nhttps://example.com/page2\nhttps://example.com/dashboard"
        url_text.insert("1.0", placeholder)
        url_text.config(fg="gray")

        def on_focus_in(event):
            if url_text.get("1.0", "end-1c") == placeholder:
                url_text.delete("1.0", "end")
                url_text.config(fg="black")

        def on_focus_out(event):
            if not url_text.get("1.0", "end-1c").strip():
                url_text.insert("1.0", placeholder)
                url_text.config(fg="gray")

        url_text.bind("<FocusIn>", on_focus_in)
        url_text.bind("<FocusOut>", on_focus_out)

        # Options frame
        options_frame = ttk.LabelFrame(main_frame, text="Options")
        options_frame.pack(fill="x", pady=(0, 10))

        wait_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(
            options_frame, text="Wait between page loads (2 seconds)", variable=wait_var
        ).pack(anchor="w", padx=5, pady=5)

        # Button frame
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill="x")

        def start_navigation():
            urls_text = url_text.get("1.0", "end-1c").strip()
            if not urls_text or urls_text == placeholder:
                messagebox.showerror("Error", "Please enter at least one URL.")
                return

            urls = [url.strip() for url in urls_text.split("\n") if url.strip()]
            if not urls:
                messagebox.showerror("Error", "Please enter valid URLs.")
                return

            dialog.destroy()
            self._navigate_and_scrape_urls(urls, wait_var.get())

        ttk.Button(button_frame, text="Cancel", command=dialog.destroy).pack(
            side="right", padx=(5, 0)
        )
        ttk.Button(
            button_frame, text="Start Navigation", command=start_navigation
        ).pack(side="right")

    def _navigate_and_scrape_urls(self, urls, wait_between_loads):
        """Navigate to URLs and scrape content using authenticated session."""
        try:
            # Show progress
            if hasattr(self.ui, "set_conversation_status"):
                self.ui.set_conversation_status("Navigating authenticated site...")
            else:
                self.ui.conversation_status_label.config(
                    text="Navigating authenticated site..."
                )
            self.update()

            # Use the navigation function
            results = navigate_and_scrape_sync(urls, wait_between_loads)

            if "error" in results:
                messagebox.showerror(
                    "Navigation Error", f"Error during navigation: {results['error']}"
                )
                return

            # Process results
            scraped_count = len(
                [r for r in results.get("results", []) if "error" not in r]
            )
            total_count = len(urls)

            # Add results to conversation
            summary = f"Navigation completed: {scraped_count}/{total_count} pages scraped successfully."
            if hasattr(self.ui, "set_conversation_status"):
                self.ui.set_conversation_status(summary)
            else:
                self.ui.conversation_status_label.config(text=summary)

            # Add to conversation
            self.conversation_history.append(
                {
                    "role": "user",
                    "content": f"Navigated and scraped {scraped_count} authenticated pages",
                }
            )

            # Update parsed files if any content was scraped
            if scraped_count > 0:
                self.parsed_files.extend(
                    [
                        {
                            "name": f"Navigation Result {i+1}: {r.get('url', 'Unknown')}",
                            "content": r.get("content", ""),
                            "url": r.get("url", ""),
                            "timestamp": r.get("timestamp", ""),
                        }
                        for i, r in enumerate(results.get("results", []))
                        if "error" not in r and r.get("content")
                    ]
                )
                self._update_parsed_file_label()

            messagebox.showinfo("Navigation Complete", summary)

        except Exception as e:
            error_msg = f"Navigation failed: {str(e)}"
            messagebox.showerror("Error", error_msg)
            if hasattr(self.ui, "set_conversation_status"):
                self.ui.set_conversation_status("Navigation failed")
            else:
                self.ui.conversation_status_label.config(text="Navigation failed")

    def run(self):
        """Run the main application loop."""
        self.mainloop()


class ToolTip:
    """Simple tooltip for Tk widgets."""

    def __init__(self, widget, text: str = "", delay_ms: int = 350):
        self.widget = widget
        self.text = text
        self.delay_ms = delay_ms
        self._tipwindow = None
        self._id = None
        self._x = self._y = 0
        widget.bind("<Enter>", self._enter, add="+")
        widget.bind("<Leave>", self._leave, add="+")
        widget.bind("<Motion>", self._motion, add="+")

    def _enter(self, _event=None):
        self._schedule()

    def _leave(self, _event=None):
        self._unschedule()
        self._hide_tip()

    def _motion(self, event):
        self._x = event.x_root + 12
        self._y = event.y_root + 8
        if self._tipwindow:
            self._tipwindow.wm_geometry(f"+{self._x}+{self._y}")

    def _schedule(self):
        self._unschedule()
        if not self.text:
            return
        self._id = self.widget.after(self.delay_ms, self._show_tip)

    def _unschedule(self):
        if self._id:
            try:
                self.widget.after_cancel(self._id)
            except Exception:
                pass
            self._id = None

    def _show_tip(self):
        if self._tipwindow or not self.text:
            return
        # Create a toplevel window
        self._tipwindow = tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{self._x}+{self._y}")
        label = ttk.Label(
            tw,
            text=self.text,
            justify="left",
            background="#ffffe0",
            relief="solid",
            borderwidth=1,
            padding=(6, 3),
        )
        label.pack(ipadx=1)

    def _hide_tip(self):
        tw = self._tipwindow
        if tw:
            tw.destroy()
            self._tipwindow = None


if __name__ == "__main__":
    app = OllamaGUI()
    app.run()
